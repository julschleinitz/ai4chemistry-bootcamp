#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Every figure in the deck, drawn as NATIVE PowerPoint objects:
  · data plots  -> real charts (right-click > Edit Data works)
  · schematics  -> real autoshapes and connectors (click and drag works)
  · tables      -> real PowerPoint tables
No raster images anywhere.
"""
import math
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import XyChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.oxml.ns import qn

from deck_lib import (PEACH, ORANGE, ORANGE_D, TEAL, TEAL_L, GREEN, PURPLE,
                      INK, GRAY, LGRAY, CARD, WHITE, FONT,
                      text, node, rbox, line, card, _margins)

MK_NONE, MK_CIRCLE, MK_TRIANGLE, MK_SQUARE = -4142, 8, 3, 1


# ===================================================================
# chart helpers
# ===================================================================
def _style_axis(ax, title=None, size=9.5, gridlines=False, tick_labels=True,
                lo=None, hi=None):
    ax.has_major_gridlines = gridlines
    if gridlines:
        gl = ax.major_gridlines.format.line
        gl.color.rgb = RGBColor(0xEC, 0xEE, 0xF2); gl.width = Pt(0.75)
    # Hide the axis line itself: on data that straddles zero it is drawn through the
    # middle of the plot and cuts across the curves. Gridlines carry the scale instead.
    ax.format.line.fill.background()
    if tick_labels:
        # LOW keeps the numbers at the edge of the plot; the default (NEXT_TO_AXIS)
        # draws them at the zero crossing, on top of the data
        ax.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        ax.tick_labels.font.size = Pt(size)
        ax.tick_labels.font.color.rgb = GRAY
        ax.tick_labels.font.name = FONT
        ax.tick_labels.number_format = "General"
        ax.tick_labels.number_format_is_linked = False
    else:
        ax.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    if lo is not None:
        ax.minimum_scale = lo
    if hi is not None:
        ax.maximum_scale = hi
    if title:
        ax.has_title = True
        ax.axis_title.text_frame.text = title
        r = ax.axis_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(size + 1); r.font.name = FONT
        r.font.color.rgb = GRAY; r.font.bold = False
    else:
        ax.has_title = False


def _pad_series(series):
    """python-pptx stacks XY series in the same spreadsheet columns. If a later
    series is longer than the first, some renderers drop the earlier series
    entirely. Padding every series to the same length with a repeat of its own
    last point is visually invisible and sidesteps the problem."""
    n = max((len(sp["pts"]) for sp in series), default=0)
    out = []
    for sp in series:
        q = dict(sp)
        pts = list(sp["pts"])
        if pts and len(pts) < n:
            pts = pts + [pts[-1]] * (n - len(pts))
        q["pts"] = pts
        out.append(q)
    return out


def xy_chart(slide, x, y, w, h, series, xlabel=None, ylabel=None,
             xlim=(None, None), ylim=(None, None), legend=False,
             gridlines=True, tick_size=9.5, tick_labels=True):
    """series: list of dicts {name, pts:[(x,y)...], kind:'markers'|'line'|'both',
    color, size, width, dash}."""
    series = _pad_series(series)
    cd = XyChartData()
    for sp in series:
        s = cd.add_series(sp.get("name", ""))
        for px, py in sp["pts"]:
            s.add_data_point(px, py)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_LINES,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.size = Pt(tick_size); ch.font.name = FONT
    for sp, ser in zip(series, ch.series):
        kind = sp.get("kind", "markers")
        col = sp.get("color", TEAL)
        if kind in ("markers", "both"):
            ser.marker.style = sp.get("marker", MK_CIRCLE)
            ser.marker.size = sp.get("size", 5)
            ser.marker.format.fill.solid()
            ser.marker.format.fill.fore_color.rgb = col
            ser.marker.format.line.color.rgb = sp.get("edge", col)
            ser.marker.format.line.width = Pt(0.75)
        else:
            ser.marker.style = MK_NONE
        if kind in ("line", "both"):
            lf = ser.format.line
            lf.color.rgb = col
            lf.width = Pt(sp.get("width", 2.25))
            if sp.get("dash"):
                lf.dash_style = sp["dash"]
        else:
            ser.format.line.fill.background()
        ser.smooth = sp.get("smooth", False)
    _style_axis(ch.value_axis, ylabel, tick_size, gridlines, tick_labels,
                ylim[0], ylim[1])
    _style_axis(ch.category_axis, xlabel, tick_size, gridlines, tick_labels,
                xlim[0], xlim[1])
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(tick_size)
        ch.legend.font.name = FONT
        ch.legend.font.color.rgb = INK
    # remove the chart's own outline so it sits cleanly on the white slide
    _no_chart_border(ch)
    return ch


def _no_chart_border(ch):
    from pptx.oxml.xmlchemy import OxmlElement as _OE
    cs = ch._chartSpace
    spPr = cs.find(qn('c:spPr'))
    if spPr is None:
        spPr = _OE('c:spPr')
        # c:spPr must follow c:chart inside c:chartSpace
        chart_el = cs.find(qn('c:chart'))
        if chart_el is not None:
            chart_el.addnext(spPr)
        else:
            cs.append(spPr)
    for tag in ('a:noFill', 'a:ln'):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    nf = _OE('a:noFill'); spPr.append(nf)
    ln = _OE('a:ln'); ln.append(_OE('a:noFill')); spPr.append(ln)


def col_chart(slide, x, y, w, h, cats, series, ylabel=None, ylim=(None, None),
              colors=None, legend=False, data_labels=False, gap=60,
              overlap=None, tick_size=9.5):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.size = Pt(tick_size); ch.font.name = FONT
    plot = ch.plots[0]
    plot.gap_width = gap
    if overlap is not None:
        plot.overlap = overlap
    if data_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(tick_size + 0.5); dl.font.bold = True
        dl.font.color.rgb = INK; dl.font.name = FONT
        dl.number_format = '0"%"'; dl.number_format_is_linked = False
    if colors:
        for ser, cols in zip(ch.series, colors):
            if isinstance(cols, list):     # per-point colouring
                for i, c in enumerate(cols):
                    pt = ser.points[i]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = c
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = cols
            ser.format.line.fill.background()
    _style_axis(ch.value_axis, ylabel, tick_size, True, True, ylim[0], ylim[1])
    _style_axis(ch.category_axis, None, tick_size, False, True)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(tick_size); ch.legend.font.name = FONT
        ch.legend.font.color.rgb = INK
    return ch


def panel_title(slide, t, x, y, w, size=12):
    text(slide, t, x, y, w, 0.28, size=size, color=INK, bold=True)


def caption(slide, t, x, y, w, size=10, color=GRAY):
    text(slide, t, x, y, w, 0.3, size=size, color=color, italic=True)


