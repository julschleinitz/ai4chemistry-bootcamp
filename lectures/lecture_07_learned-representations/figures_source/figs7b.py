#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Figures for lecture_07, parts 3-6: autoencoders & VAEs, self-supervised
pre-training, transfer learning, and the failure modes.

Same rule as figs7.py — native PowerPoint shapes, charts and tables only.
"""
import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH

from deck_lib import (PEACH, ORANGE, ORANGE_D, TEAL, TEAL_L, GREEN, PURPLE,
                      INK, GRAY, LGRAY, CARD, WHITE, FONT, SW, SH,
                      BODY_TOP, BODY_BOT,
                      text, node, rbox, line, card, centered, table, _margins)
from chartlib import xy_chart, col_chart, panel_title, caption
from figs7 import (BLUSH, MINT, SAND, RED, OLIVE, notes, arrow, chevron, blob,
                   dot, _rng, scatter_cloud, mol_sketch, stack, legend, ref)


# =====================================================================
# PART 3 — autoencoders and VAEs
# =====================================================================
def _trapezoid(slide, x, y, w, h, flip=False, fill=TEAL_L, label="", size=11.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.rotation = 270 if not flip else 90
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    if label:
        tf = sh.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
    return sh


def fig_autoencoder(s, y=1.4, h=4.05):
    """encoder -> bottleneck -> decoder, with the reconstruction loss closing the loop."""
    cy = y + 1.55
    # input molecule
    card(s, 0.62, y + 0.55, 1.85, 2.05, fill=CARD)
    mol_sketch(s, 1.54, cy - 0.10, sc=1.35, tails=3, seed=4, color=INK, hub=ORANGE)
    text(s, "input  x", 0.62, y + 2.22, 1.85, 0.3, size=11.5, color=INK,
         bold=True, align=PP_ALIGN.CENTER)
    # encoder
    _trapezoid(s, 2.92, y + 0.62, 1.90, 1.55, flip=False, fill=TEAL, label="")
    text(s, "ENCODER", 2.82, cy - 0.16, 2.10, 0.3, size=12, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    # bottleneck
    bx = 6.05
    for k in range(6):
        rbox(s, bx, y + 0.92 + k * 0.24, 0.60, 0.19, fill=ORANGE, edge=ORANGE,
             label="", size=8)
    card(s, bx - 0.30, y + 0.72, 1.20, 1.86, fill=BLUSH)
    for k in range(6):
        rbox(s, bx, y + 0.92 + k * 0.24, 0.60, 0.19, fill=ORANGE, edge=ORANGE,
             label="", size=8)
    text(s, "z", bx - 0.30, y + 2.62, 1.20, 0.3, size=15, color=ORANGE_D,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "the bottleneck\n196 numbers", bx - 0.95, y + 2.94, 2.50, 0.55,
         size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    # decoder
    _trapezoid(s, 7.98, y + 0.62, 1.90, 1.55, flip=True, fill=TEAL_L, label="")
    text(s, "DECODER", 7.88, cy - 0.16, 2.10, 0.3, size=12, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    # output
    card(s, 11.05, y + 0.55, 1.85, 2.05, fill=CARD)
    mol_sketch(s, 11.97, cy - 0.10, sc=1.35, tails=3, seed=4, color=TEAL)
    text(s, "output  x̂", 11.05, y + 2.22, 1.85, 0.3, size=11.5, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    for (a, b) in [(2.52, 2.88), (4.88, 6.00), (6.92, 7.94), (9.92, 11.00)]:
        arrow(s, a, cy, b, cy, color=INK, w=1.8)
    # loss loop, routed below the bottleneck caption
    ly = y + 3.72
    line(s, 1.54, y + 2.62, 1.54, ly, color=RED, w=1.6, dash=DASH.DASH)
    line(s, 11.97, y + 2.62, 11.97, ly, color=RED, w=1.6, dash=DASH.DASH)
    line(s, 1.54, ly, 11.97, ly, color=RED, w=1.6, dash=DASH.DASH)
    card(s, 5.30, ly - 0.24, 2.85, 0.50, fill=WHITE)
    text(s, "the loss is:  make x̂ = x", 5.30, ly - 0.17, 2.85, 0.36, size=12.5,
         color=RED, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_bottleneck_why(s, y=1.48, h=3.75):
    """Why a bottleneck produces a useful representation."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    panels = [
        ("bottleneck too wide", "z is as big as x",
         "the network copies the input\nand learns nothing", RED, BLUSH),
        ("bottleneck just right", "z ≪ x",
         "it must keep ring systems,\nfunctional groups, size —\nand throw away the rest", GREEN, MINT),
        ("bottleneck too narrow", "z ≈ 2",
         "everything reconstructs\nto benzene", RED, BLUSH),
    ]
    for x, (hd, sub, body, col, tint) in zip(xs, panels):
        card(s, x, y, W, h, fill=tint)
        text(s, hd, x, y + 0.16, W, 0.3, size=12.5, color=col, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, sub, x, y + 0.48, W, 0.3, size=11, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)
    # widths cartoon
    for i, (x, nb) in enumerate(zip(xs, [11, 5, 1])):
        cy = y + 1.68
        for k in range(nb):
            rbox(s, x + W / 2 - 0.32, cy - (nb * 0.155) / 2 + k * 0.155, 0.64, 0.12,
                 fill=ORANGE, edge=ORANGE, label="", size=8)
    for x, (hd, sub, body, col, tint) in zip(xs, panels):
        text(s, body, x + 0.22, y + h - 1.12, W - 0.44, 1.0, size=11.5, color=INK,
             align=PP_ALIGN.CENTER)
    card(s, 0.62, y + h + 0.14, SW - 1.24, 0.52, fill=PEACH)
    text(s, "compression is the whole trick: what survives the bottleneck is what the "
            "data says matters", 0.62, y + h + 0.22, SW - 1.24, 0.38, size=13.5,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    return y + h + 0.7


def fig_ae_holes(s, y=1.42, h=3.95):
    """Plain AE: the latent space is a set of islands with garbage in between."""
    card(s, 0.62, y, 5.85, h, fill=CARD)
    panel_title(s, "plain autoencoder — latent space", 0.85, y + 0.13, 5.4, size=11.5)
    r = _rng(17)
    islands = [(2.05, y + 1.35), (4.30, y + 1.20), (3.05, y + 2.55),
               (5.10, y + 2.70), (1.60, y + 2.85)]
    for k, (ix, iy) in enumerate(islands):
        scatter_cloud(s, ix, iy, 0.40, 0.32, 14, 60 + k, fill=TEAL, d=0.10)
    # sample a hole
    hx, hy = 3.30, y + 1.90
    dot(s, hx, hy, d=0.20, fill=RED)
    text(s, "★", hx - 0.16, hy - 0.42, 0.35, 0.3, size=14, color=RED, bold=True)
    text(s, "sample here…", hx + 0.18, hy - 0.14, 1.5, 0.28, size=10.5, color=RED,
         bold=True)
    text(s, "the model was never asked what lives between the islands",
         0.85, y + h - 0.44, 5.4, 0.3, size=10.5, color=GRAY, italic=True,
         align=PP_ALIGN.CENTER)
    arrow(s, 6.62, y + 1.70, 7.40, y + 1.70, color=RED, w=2.0)
    text(s, "decode", 6.55, y + 1.32, 0.95, 0.28, size=10.5, color=RED,
         align=PP_ALIGN.CENTER)
    # right: garbage output
    card(s, 7.60, y, 5.12, h, fill=BLUSH)
    panel_title(s, "what comes out", 7.85, y + 0.13, 4.6, size=11.5)
    bad = ["C1=CC=CC=C1C(=O)N(C)(C)(C)C", "c1ccc2cc(", "CC(C)(C)(C)(C)C",
           "N#N#N=O", "C1CC1CC1CC1CC…"]
    for k, t in enumerate(bad):
        rbox(s, 7.90, y + 0.62 + k * 0.50, 4.55, 0.40, fill=WHITE, edge=LGRAY,
             label=t, size=10.5, label_color=INK, bold=False)
        text(s, "✗", 12.28, y + 0.66 + k * 0.50, 0.3, 0.3, size=13, color=RED, bold=True)
    text(s, "invalid, unparseable, or chemically absurd", 7.85, y + h - 0.44, 4.6, 0.3,
         size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_vae_point_vs_blob(s, y=1.45, h=3.85):
    """The one idea of a VAE: encode a region, not a point."""
    for i, (x, hd, tint) in enumerate([(0.62, "autoencoder: x → one point", CARD),
                                       (6.86, "VAE: x → a small cloud", MINT)]):
        card(s, x, y, 5.86, h, fill=tint)
        text(s, hd, x, y + 0.14, 5.86, 0.3, size=12.5,
             color=GRAY if i == 0 else TEAL, bold=True, align=PP_ALIGN.CENTER)
        mol_sketch(s, x + 0.95, y + 1.75, sc=1.15, tails=2, seed=4, color=INK)
        arrow(s, x + 1.60, y + 1.75, x + 2.35, y + 1.75, color=INK, w=1.8)
        cx, cy = x + 3.85, y + 1.75
        if i == 0:
            dot(s, cx, cy, d=0.22, fill=ORANGE)
            text(s, "z", cx + 0.16, cy - 0.36, 0.4, 0.3, size=13, color=ORANGE_D, bold=True)
        else:
            for (rx, ry, col) in [(0.78, 0.62, RGBColor(0xEC, 0xF4, 0xF7)),
                                  (0.52, 0.41, RGBColor(0xD2, 0xE6, 0xEF)),
                                  (0.26, 0.21, TEAL_L)]:
                blob(s, cx, cy, rx, ry, fill=col)
            dot(s, cx, cy, d=0.16, fill=ORANGE)
            for (ax, ay) in [(-0.42, 0.22), (0.35, -0.28), (0.12, 0.40), (-0.30, -0.34)]:
                dot(s, cx + ax, cy + ay, d=0.11, fill=ORANGE_D)
            text(s, "mean  +  spread", cx - 1.15, cy + 0.85, 2.3, 0.3, size=11,
                 color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "training on samples from the cloud forces nearby points to decode to "
            "similar molecules — the gaps get filled",
         0.62, y + h + 0.16, SW - 1.24, 0.3, size=12.5, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h + 0.5


def fig_vae_forces(s, y=1.42, h=3.95):
    """Reconstruction vs regularisation, as a tug of war."""
    cy = y + 1.62
    # left force
    card(s, 0.62, y, 3.95, h, fill=BLUSH)
    text(s, "RECONSTRUCTION", 0.62, y + 0.16, 3.95, 0.3, size=12.5, color=ORANGE_D,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "“keep every molecule distinguishable”", 0.72, y + 0.48, 3.75, 0.3,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    for k, (ix, iy) in enumerate([(1.35, cy - 0.45), (3.55, cy - 0.30),
                                  (1.60, cy + 0.75), (3.30, cy + 0.85)]):
        scatter_cloud(s, ix, iy, 0.30, 0.24, 7, 90 + k, fill=ORANGE, d=0.10)
    text(s, "pushes clouds APART", 0.62, y + h - 0.66, 3.95, 0.3, size=11.5,
         color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    text(s, "→ islands and holes come back", 0.62, y + h - 0.38, 3.95, 0.3,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    # centre
    card(s, 4.75, y + 0.95, 3.85, 1.75, fill=WHITE)
    text(s, "the VAE loss", 4.75, y + 1.10, 3.85, 0.3, size=12.5, color=INK,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "reconstruct  +  β × stay-organised", 4.75, y + 1.48, 3.85, 0.32,
         size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "β is a dial you turn", 4.75, y + 1.86, 3.85, 0.3, size=11,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    text(s, "β too small → autoencoder    β too large → mush", 4.75, y + 2.16,
         3.85, 0.3, size=10.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    arrow(s, 4.66, cy, 4.20, cy, color=ORANGE_D, w=2.4)
    arrow(s, 8.68, cy, 9.14, cy, color=TEAL, w=2.4)
    # right force
    card(s, 8.77, y, 3.95, h, fill=MINT)
    text(s, "REGULARISATION", 8.77, y + 0.16, 3.95, 0.3, size=12.5, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "“all clouds should look like one blob”", 8.87, y + 0.48, 3.75, 0.3,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    blob(s, 10.75, cy + 0.22, 1.28, 1.00, fill=WHITE, edge=TEAL)
    scatter_cloud(s, 10.75, cy + 0.22, 1.00, 0.78, 34, 99, fill=TEAL, d=0.10)
    text(s, "pulls clouds TOGETHER", 8.77, y + h - 0.66, 3.95, 0.3, size=11.5,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "→ one smooth space you can sample", 8.77, y + h - 0.38, 3.95, 0.3,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_interpolation(s, y=1.45, h=3.9):
    """Walk a straight line in latent space, decode at every step."""
    card(s, 0.62, y, SW - 1.24, h, fill=CARD)
    panel_title(s, "a straight line in z is a smooth series of molecules in chemistry",
                0.88, y + 0.14, 8.0, size=11.5)
    cy = y + 1.65
    n = 6
    x0, x1 = 1.35, 11.95
    for k in range(n):
        px = x0 + k * (x1 - x0) / (n - 1)
        end = (k == 0 or k == n - 1)
        col = ORANGE if end else TEAL
        blob(s, px, cy, 0.52, 0.52, fill=BLUSH if end else MINT)
        mol_sketch(s, px, cy, sc=1.0, tails=2 if k < 3 else 3, seed=4 + k, color=INK,
                   hub=col if end else None)
        if k < n - 1:
            arrow(s, px + 0.60, cy, px + (x1 - x0) / (n - 1) - 0.60, cy,
                  color=LGRAY, w=1.6)
    text(s, "start molecule", x0 - 0.85, cy + 0.78, 1.7, 0.28, size=10.5,
         color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    text(s, "target molecule", x1 - 0.85, cy + 0.78, 1.7, 0.28, size=10.5,
         color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    line(s, x0, cy + 1.20, x1, cy + 1.20, color=INK, w=1.6)
    for k in range(n):
        px = x0 + k * (x1 - x0) / (n - 1)
        dot(s, px, cy + 1.20, d=0.13, fill=INK)
    text(s, "z₀", x0 - 0.3, cy + 1.32, 0.6, 0.26, size=10.5, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, "z₁", x1 - 0.3, cy + 1.32, 0.6, 0.26, size=10.5, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, "this is what people mean by “a continuous representation of molecules”",
         0.88, y + h - 0.44, SW - 1.76, 0.3, size=12, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_chemvae(s, y=1.35, h=4.15):
    """ChemVAE: SMILES in, SMILES out, with a property head bending the latent space."""
    cy = y + 1.30
    rbox(s, 0.62, cy - 0.32, 1.85, 0.66, fill=WHITE, edge=LGRAY,
         label="SMILES\nCC(=O)Oc1ccccc1", size=9.5, label_color=INK, bold=False)
    _trapezoid(s, 2.85, cy - 0.75, 1.55, 1.35, fill=TEAL, label="")
    text(s, "encoder", 2.80, cy - 0.13, 1.65, 0.3, size=11, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "CNN", 2.80, cy - 0.98, 1.65, 0.26, size=10, color=GRAY,
         align=PP_ALIGN.CENTER)
    card(s, 5.05, cy - 0.78, 1.30, 1.55, fill=BLUSH)
    for k in range(5):
        rbox(s, 5.35, cy - 0.58 + k * 0.26, 0.70, 0.20, fill=ORANGE, edge=ORANGE,
             label="", size=8)
    text(s, "z  ·  196-D", 4.75, cy + 0.84, 1.90, 0.3, size=11, color=ORANGE_D,
         bold=True, align=PP_ALIGN.CENTER)
    _trapezoid(s, 6.95, cy - 0.75, 1.55, 1.35, flip=True, fill=TEAL_L, label="")
    text(s, "decoder", 6.90, cy - 0.13, 1.65, 0.3, size=11, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "GRU", 6.90, cy - 0.98, 1.65, 0.26, size=10, color=GRAY,
         align=PP_ALIGN.CENTER)
    rbox(s, 9.00, cy - 0.32, 1.85, 0.66, fill=WHITE, edge=TEAL,
         label="SMILES\nCC(=O)Oc1ccccc1", size=9.5, label_color=INK, bold=False)
    for (a, b) in [(2.52, 2.80), (4.45, 5.00), (6.40, 6.90), (8.55, 8.95)]:
        arrow(s, a, cy, b, cy, color=INK, w=1.7)
    # property head
    arrow(s, 5.70, cy + 0.86, 5.70, cy + 1.30, color=PURPLE, w=2.0)
    rbox(s, 4.60, cy + 1.32, 2.20, 0.58, fill=WHITE, edge=PURPLE,
         label="property head\nlogP, QED, SAS", size=10, label_color=PURPLE)
    text(s, "trained jointly — this is\nwhat organises the space", 7.05, cy + 1.30,
         3.3, 0.6, size=11, color=PURPLE, bold=True)
    # numbers strip
    facts = [("250,000", "ZINC molecules, no labels"),
             ("196", "latent dimensions"),
             ("~4 %", "of random z decode to a valid molecule")]
    bw = 3.86
    for i, (big, small) in enumerate(facts):
        x = 0.62 + i * (bw + 0.32)
        card(s, x, y + 3.18, bw, 0.96, fill=CARD)
        text(s, big, x, y + 3.26, bw, 0.40, size=21,
             color=RED if i == 2 else TEAL, bold=True, align=PP_ALIGN.CENTER)
        text(s, small, x + 0.14, y + 3.68, bw - 0.28, 0.36, size=10.5, color=GRAY,
             align=PP_ALIGN.CENTER)
    return y + h


def fig_validity_race(s, y=1.45, h=3.95):
    """The validity arms race — and the note that the last two are by construction."""
    col_chart(s, 1.35, y, 10.65, h - 0.72,
              ["character VAE\n(ChemVAE)", "grammar VAE", "syntax-directed\nVAE",
               "junction-tree VAE", "SELFIES\n(any model)"],
              [("% of random latent points that decode to a valid molecule",
                [0.7, 7.2, 43.5, 100.0, 100.0])],
              ylabel="prior validity  (%)", ylim=(0, 140),
              colors=[[RED, ORANGE, ORANGE_D, TEAL, TEAL]],
              data_labels=True, gap=62, tick_size=10.5)
    x0 = 1.35 + 0.80
    wcol = (10.65 - 0.80) / 5
    line(s, x0 + 3 * wcol, y + 0.46, x0 + 5 * wcol - 0.15, y + 0.46, color=TEAL, w=2.0)
    text(s, "by construction, not by learning", x0 + 3 * wcol - 0.35,
         y + 0.10, 2 * wcol + 0.55, 0.3, size=12, color=TEAL, bold=True,
         align=PP_ALIGN.CENTER)
    line(s, x0 + 0.10, y + 0.46, x0 + 3 * wcol - 0.30, y + 0.46, color=GRAY, w=1.4)
    text(s, "learned", x0 + 0.10, y + 0.10, 3 * wcol - 0.40, 0.3, size=12,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_jtvae_selfies(s, y=1.4, h=4.05):
    """Two ways to make invalidity impossible."""
    # left: junction tree
    card(s, 0.62, y, 5.90, h, fill=MINT)
    text(s, "JT-VAE — assemble from a vocabulary of valid pieces", 0.80, y + 0.14,
         5.55, 0.3, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    mol_sketch(s, 1.75, y + 1.30, sc=1.35, tails=2, seed=4, color=INK)
    arrow(s, 2.62, y + 1.30, 3.20, y + 1.30, color=INK, w=1.8)
    tp = {0: (4.15, y + 0.92), 1: (3.55, y + 1.72), 2: (4.85, y + 1.75),
          3: (4.20, y + 2.48)}
    for a, b in [(0, 1), (0, 2), (1, 3)]:
        line(s, tp[a][0], tp[a][1], tp[b][0], tp[b][1], color=LGRAY, w=2.2)
    for k, (px, py) in tp.items():
        node(s, px, py, 0.44, fill=WHITE, edge=TEAL, label=["ring", "C=O", "ring", "N"][k],
             size=9, label_color=TEAL, edge_w=1.6)
    text(s, "junction tree of substructures", 3.05, y + 2.86, 2.6, 0.28, size=10,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    text(s, "780-piece vocabulary · illegal joins are masked out during decoding",
         0.80, y + h - 0.72, 5.55, 0.55, size=11, color=INK, align=PP_ALIGN.CENTER)
    text(s, "100 % valid", 0.80, y + h - 0.38, 5.55, 0.3, size=12.5, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    # right: SELFIES
    card(s, 6.82, y, 5.90, h, fill=MINT)
    text(s, "SELFIES — a grammar that cannot fail", 7.00,
         y + 0.14, 5.55, 0.3, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    toks = ["[F]", "[=C]", "[=C]", "[#N]"]
    states = ["X₀", "X₁", "X₃", "X₂"]
    outs = ["F", "F–C", "F–C=C", "F–C=C=N"]
    tw = 1.24
    for k in range(4):
        x = 7.10 + k * (tw + 0.16)
        rbox(s, x, y + 0.62, tw, 0.44, fill=WHITE, edge=TEAL, label=toks[k],
             size=11, label_color=INK)
        text(s, states[k], x, y + 1.10, tw, 0.26, size=9.5, color=GRAY,
             align=PP_ALIGN.CENTER)
        rbox(s, x, y + 1.40, tw, 0.44, fill=BLUSH, edge=BLUSH, label=outs[k],
             size=10, label_color=INK, bold=False)
        if k < 3:
            arrow(s, x + tw + 0.02, y + 0.84, x + tw + 0.14, y + 0.84, color=GRAY, w=1.4)
    text(s, "asked for a triple bond, valence allowed only a double →  clamped, "
            "never rejected", 7.00, y + 2.02, 5.55, 0.55, size=10.5, color=ORANGE_D,
         italic=True, align=PP_ALIGN.CENTER)
    text(s, "100 % valid  ·  but only local valence is enforced —\n"
            "not aromaticity, not ring strain, not stability",
         7.00, y + h - 0.90, 5.55, 0.6, size=11, color=INK, align=PP_ALIGN.CENTER)
    text(s, "1 random mutation:  SMILES 9.9 %  →  SELFIES 100 % valid", 7.00,
         y + h - 0.38, 5.55, 0.3, size=11.5, color=TEAL, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_honest_turn(s, y=1.38, h=4.1):
    """Validity was the easy metric — three inconvenient numbers."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    # 1 — memorisation
    x = xs[0]
    card(s, x, y, W, h, fill=BLUSH)
    text(s, "it memorises", x, y + 0.16, W, 0.3, size=13, color=RED, bold=True,
         align=PP_ALIGN.CENTER)
    col_chart(s, x + 0.26, y + 0.58, W - 0.52, 1.95,
              ["char-RNN", "VAE"], [("novelty", [0.842, 0.695])],
              ylabel="novelty (MOSES)", ylim=(0, 1.0), colors=[[TEAL, RED]],
              gap=80, tick_size=9.5)
    text(s, "31 % of what the VAE “designs”\nis already in the training set",
         x + 0.20, y + 2.68, W - 0.40, 0.62, size=11, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, "MOSES benchmark", x, y + h - 0.36, W, 0.3, size=10, color=GRAY,
         italic=True, align=PP_ALIGN.CENTER)
    # 2 — validity is not quality
    x = xs[1]
    card(s, x, y, W, h, fill=BLUSH)
    text(s, "100 % valid ≠ good", x, y + 0.16, W, 0.3, size=13, color=RED,
         bold=True, align=PP_ALIGN.CENTER)
    rows = [["Molecule Chef", "99.0", "0.73"], ["CGVAE", "100.0", "11.73"],
            ["char-VAE", "12.0", "37.65"]]
    table(s, ["model", "valid %", "FCD ↓"], rows, x=x + 0.20, y=y + 0.62,
          w=W - 0.40, size=10.5, head_size=10.5, row_h=1.72)
    text(s, "the model with perfect validity\nmatches the real distribution worst",
         x + 0.20, y + 2.52, W - 0.40, 0.62, size=11, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, "FCD = distance to real chemistry", x, y + h - 0.36, W, 0.3, size=10,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    # 3 — a GA wins
    x = xs[2]
    card(s, x, y, W, h, fill=BLUSH)
    text(s, "a genetic algorithm wins", x, y + 0.16, W, 0.3, size=13, color=RED,
         bold=True, align=PP_ALIGN.CENTER)
    col_chart(s, x + 0.26, y + 0.58, W - 0.52, 1.95,
              ["graph GA\n30 s", "ChemTS\n8 h", "char-VAE\n+ BO, 8 h"],
              [("score", [7.4, 5.6, 0.0])],
              ylabel="best penalised logP", ylim=(0, 8.5),
              colors=[[GREEN, TEAL, RED]], gap=55, tick_size=9.5)
    text(s, "30 seconds on a laptop beats\n8 hours of VAE + Bayesian opt.",
         x + 0.20, y + 2.68, W - 0.40, 0.62, size=11, color=INK,
         align=PP_ALIGN.CENTER)
    text(s, "Jensen, Chem. Sci. 2019", x, y + h - 0.36, W, 0.3, size=10,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_vae_verdict(s, y=1.48, h=3.8):
    """What a VAE is actually good for."""
    card(s, 0.62, y, 5.95, h, fill=BLUSH)
    text(s, "as a molecule generator", 0.62, y + 0.18, 5.95, 0.32, size=14,
         color=RED, bold=True, align=PP_ALIGN.CENTER)
    items = ["beaten by a graph genetic algorithm",
             "beaten by an ordinary SMILES LSTM",
             "beaten by random screening of ZINC",
             "only ~30 % of proposals are synthesisable"]
    for k, t in enumerate(items):
        text(s, "✗", 1.05, y + 0.72 + k * 0.52, 0.3, 0.3, size=14, color=RED, bold=True)
        text(s, t, 1.42, y + 0.74 + k * 0.52, 4.9, 0.4, size=12, color=INK)
    text(s, "the field moved to diffusion, GAs and LLMs\n(→ Thursday's lecture)",
         0.85, y + h - 0.92, 5.5, 0.6, size=11, color=GRAY, italic=True,
         align=PP_ALIGN.CENTER)
    card(s, 6.77, y, 5.95, h, fill=MINT)
    text(s, "as a representation learner", 6.77, y + 0.18, 5.95, 0.32, size=14,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    items = ["a fixed-length vector for any molecule",
             "trained without a single label",
             "a smooth space you can optimise inside",
             "the encoder is reusable — that is the part\nthat survived"]
    for k, t in enumerate(items):
        text(s, "✓", 7.20, y + 0.72 + k * 0.52, 0.3, 0.3, size=14, color=GREEN, bold=True)
        text(s, t, 7.57, y + 0.74 + k * 0.52, 4.9, 0.5, size=12, color=INK)
    text(s, "keep the encoder, throw away the decoder", 7.00, y + h - 0.62, 5.5, 0.32,
         size=12.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    return y + h


# =====================================================================
# PART 4 — self-supervised pre-training
# =====================================================================
def fig_two_families(s, y=1.4, h=4.05):
    """Masked prediction vs contrastive learning."""
    # left: masked prediction
    card(s, 0.62, y, 5.95, h, fill=MINT)
    text(s, "A · MASKED PREDICTION", 0.62, y + 0.16, 5.95, 0.3, size=13,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "hide part of the molecule, predict it back", 0.62, y + 0.48, 5.95, 0.3,
         size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    toks = ["C", "C", "(", "=", "?", ")", "O", "c", "1"]
    tw = 0.54
    bx = 0.62 + (5.95 - len(toks) * (tw + 0.06)) / 2
    for k, tk in enumerate(toks):
        hid = (tk == "?")
        rbox(s, bx + k * (tw + 0.06), y + 1.02, tw, 0.50,
             fill=PURPLE if hid else WHITE, edge=PURPLE if hid else LGRAY,
             label=tk, size=12, label_color=WHITE if hid else INK)
    arrow(s, 3.60, y + 1.62, 3.60, y + 2.05, color=PURPLE, w=2.0)
    rbox(s, 2.35, y + 2.08, 2.50, 0.52, fill=WHITE, edge=PURPLE,
         label="it was “O”", size=12, label_color=PURPLE)
    text(s, "to fill the gap the model must learn valence,\nring closure and "
            "functional-group grammar",
         0.90, y + h - 0.98, 5.4, 0.6, size=11, color=INK, align=PP_ALIGN.CENTER)
    text(s, "ChemBERTa · MoLFormer · attribute masking on graphs", 0.90, y + h - 0.38,
         5.4, 0.3, size=10.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    # right: contrastive
    card(s, 6.77, y, 5.95, h, fill=MINT)
    text(s, "B · CONTRASTIVE", 6.77, y + 0.16, 5.95, 0.3, size=13, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "two damaged views of one molecule must agree", 6.77, y + 0.48, 5.95, 0.3,
         size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    # two augmented views of one molecule (left column) ...
    for cy2, sd, tl, lbl in [(y + 1.28, 4, 2, "view 1"), (y + 2.46, 9, 1, "view 2")]:
        blob(s, 8.22, cy2, 0.50, 0.44, fill=WHITE, edge=TEAL)
        mol_sketch(s, 8.22, cy2, sc=0.85, tails=tl, seed=sd, color=INK)
        text(s, lbl, 6.92, cy2 - 0.13, 0.76, 0.26, size=9.5, color=GRAY,
             align=PP_ALIGN.RIGHT)
    # ... pulled together
    arrow(s, 8.94, y + 1.50, 8.94, y + 2.26, color=GREEN, w=2.4)
    arrow(s, 8.94, y + 2.26, 8.94, y + 1.50, color=GREEN, w=2.4)
    text(s, "pull\ntogether", 9.04, y + 1.70, 1.0, 0.5, size=10.5, color=GREEN,
         bold=True)
    # ... and pushed away from a different molecule (right column)
    blob(s, 11.78, y + 1.86, 0.50, 0.44, fill=WHITE, edge=LGRAY)
    mol_sketch(s, 11.78, y + 1.86, sc=0.85, tails=3, seed=15, color=GRAY)
    text(s, "a different\nmolecule", 11.18, y + 2.42, 1.2, 0.48, size=9.5,
         color=GRAY, align=PP_ALIGN.CENTER)
    arrow(s, 10.16, y + 1.86, 11.20, y + 1.86, color=RED, w=2.2)
    text(s, "push apart", 9.98, y + 1.48, 1.4, 0.3, size=10.5, color=RED, bold=True)
    text(s, "MolCLR · GraphCL · InfoMax", 6.95, y + h - 0.38, 5.4, 0.3, size=10.5,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_molclr_augs(s, y=1.45, h=3.8):
    """The three graph augmentations, and the honest ablation result."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    heads = ["atom masking", "bond deletion", "subgraph removal"]
    subs = ["replace an atom type\nwith a dummy",
            "drop a bond,\nkeep the atoms",
            "delete a whole\nconnected region"]
    for i, (x, hd, sub) in enumerate(zip(xs, heads, subs)):
        card(s, x, y, W, h, fill=CARD)
        text(s, hd, x, y + 0.16, W, 0.3, size=13, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, sub, x, y + 0.48, W, 0.55, size=10.5, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)
        cx, cy = x + W / 2, y + 1.95
        P = {0: (cx, cy - 0.55), 1: (cx - 0.62, cy - 0.18), 2: (cx + 0.62, cy - 0.18),
             3: (cx - 0.42, cy + 0.52), 4: (cx + 0.42, cy + 0.52),
             5: (cx - 1.20, cy + 0.10), 6: (cx + 1.20, cy + 0.14)}
        E = [(0, 1), (0, 2), (1, 3), (2, 4), (3, 4), (1, 5), (2, 6)]
        killed_e, killed_n, masked_n = set(), set(), set()
        if i == 0:
            masked_n = {5}
        elif i == 1:
            killed_e = {(3, 4)}
        else:
            killed_n = {2, 4, 6}
            killed_e = {(0, 2), (2, 4), (3, 4), (2, 6)}
        for a, b in E:
            col = LGRAY if (a, b) not in killed_e else RED
            dash = None if (a, b) not in killed_e else DASH.DASH
            line(s, P[a][0], P[a][1], P[b][0], P[b][1], color=col, w=2.0, dash=dash)
        for k, (px, py) in P.items():
            if k in masked_n:
                node(s, px, py, 0.34, fill=RED, edge=RED, label="?", size=10,
                     label_color=WHITE, edge_w=1.4)
            elif k in killed_n:
                node(s, px, py, 0.30, fill=WHITE, edge=RED, label="", edge_w=1.4)
            else:
                node(s, px, py, 0.30, fill=TEAL, edge=TEAL, label="", edge_w=1.4)
    text(s, "ablation: subgraph removal alone works best — combining all three "
            "destroys too much topology, and on BBBP it makes things worse",
         0.62, y + h + 0.14, SW - 1.24, 0.32, size=12, color=RED, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h + 0.5


def fig_corpus_scaling(s, y=1.4, h=4.05):
    """Bigger corpora — and the two places bigger did not help."""
    xy_chart(s, 0.80, y + 0.30, 6.10, h - 0.95,
             [{"name": "BBBP ROC-AUC", "pts": [(6.0, 64.3), (7.0, 90.9), (8.0, 91.5),
                                               (9.04, 93.7)],
               "kind": "both", "color": TEAL, "size": 8, "width": 2.6}],
             xlabel="pre-training corpus  (log₁₀ molecules)",
             ylabel="BBBP ROC-AUC  (%)", xlim=(5.5, 9.6), ylim=(60, 100))
    for lx, ly, lbl, dx, dy in [(6.0, 64.3, "ChemBERTa\n10 M", 0.98, -0.24),
                                (9.04, 93.7, "MoLFormer-XL\n1.1 B", -0.30, -0.68)]:
        px = 0.80 + 0.80 + (lx - 5.5) / 4.1 * (6.10 - 0.95)
        py = y + 0.30 + (1 - (ly - 60) / 40) * (h - 0.95 - 0.55)
        text(s, lbl, px - 1.0 + dx, py + dy, 2.0, 0.5,
             size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "scale helps…", 0.80, y + 0.02, 6.10, 0.3, size=13, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "1.1 B SMILES on 16 GPUs — linear attention made it affordable",
         0.80, y + h - 0.50, 6.10, 0.3, size=10.5, color=GRAY, italic=True,
         align=PP_ALIGN.CENTER)
    # right: where it backfired
    card(s, 7.20, y, 5.52, h, fill=BLUSH)
    text(s, "…until it doesn't", 7.20, y + 0.10, 5.52, 0.3, size=13, color=RED,
         bold=True, align=PP_ALIGN.CENTER)
    col_chart(s, 7.50, y + 0.52, 4.95, h - 1.40,
              ["5 M", "10 M", "77 M"],
              [("ClinTox ROC-AUC", [0.341, 0.349, 0.239])],
              ylabel="ClinTox ROC-AUC", ylim=(0, 0.45),
              colors=[[ORANGE, ORANGE, RED]], gap=70, tick_size=10)
    text(s, "ChemBERTa-2, masked-LM: a 15× bigger corpus made this task worse",
         7.42, y + h - 0.72, 5.10, 0.55, size=11, color=INK, align=PP_ALIGN.CENTER)
    text(s, "corpus composition matters more than corpus size", 7.42, y + h - 0.36,
         5.10, 0.3, size=11.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_hu_pretraining(s, y=1.42, h=3.95):
    """Node-level and graph-level pretraining have to be combined."""
    ch = col_chart(s, 0.85, y + 0.20, 7.20, h - 0.70,
                   ["no\npre-training", "node-level\nonly", "graph-level\nonly",
                    "node + graph"],
                   [("average ROC-AUC", [67.0, 71.1, 70.0, 74.2])],
                   ylabel="average ROC-AUC over 8 datasets  (%)", ylim=(60, 78),
                   colors=[[GRAY, TEAL_L, TEAL_L, TEAL]], data_labels=True,
                   gap=62, tick_size=10.5)
    ch.plots[0].data_labels.number_format = '0.0'
    text(s, "+7.2 over no pre-training", 5.10, y + 0.18, 3.0, 0.32, size=13,
         color=GREEN, bold=True)
    card(s, 8.55, y + 0.20, 4.17, h - 0.70, fill=BLUSH)
    text(s, "the same recipe, different backbone", 8.70, y + 0.36, 3.87, 0.3,
         size=11.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    gains = [("GIN", 7.2, GREEN), ("GCN", 3.4, GREEN), ("GraphSAGE", 2.0, GREEN),
             ("GAT", -6.5, RED)]
    for k, (nm, g, col) in enumerate(gains):
        yy = y + 0.86 + k * 0.52
        text(s, nm, 8.80, yy, 1.6, 0.3, size=12, color=INK, bold=True)
        w = abs(g) / 8.0 * 1.45
        x0 = 10.75
        if g > 0:
            rbox(s, x0, yy + 0.04, w, 0.24, fill=col, edge=col, label="")
        else:
            rbox(s, x0 - w, yy + 0.04, w, 0.24, fill=col, edge=col, label="")
        text(s, f"{g:+.1f}", 12.05, yy, 0.6, 0.3, size=11.5, color=col, bold=True)
    line(s, 10.75, y + 0.82, 10.75, y + 0.82 + 4 * 0.52, color=GRAY, w=1.0)
    text(s, "pre-training a weaker architecture actively hurts",
         8.70, y + h - 0.86, 3.87, 0.55, size=11, color=INK, align=PP_ALIGN.CENTER)
    return y + h


# =====================================================================
# PART 5 — transfer learning
# =====================================================================
def fig_pretrain_finetune(s, y=1.38, h=4.1):
    """The two-stage recipe, drawn as one wide funnel."""
    # stage 1
    card(s, 0.62, y, 6.05, h, fill=MINT)
    text(s, "STAGE 1 · PRE-TRAIN", 0.62, y + 0.16, 6.05, 0.3, size=13, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "once, by someone with a GPU cluster", 0.62, y + 0.48, 6.05, 0.3,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    blob(s, 2.05, y + 1.55, 1.10, 0.72, fill=WHITE, edge=TEAL)
    scatter_cloud(s, 2.05, y + 1.55, 0.88, 0.55, 46, 3, fill=TEAL, d=0.09)
    text(s, "10⁶ – 10⁹ molecules\nno labels", 1.10, y + 2.38, 1.95, 0.55, size=10.5,
         color=INK, align=PP_ALIGN.CENTER)
    arrow(s, 3.30, y + 1.55, 3.90, y + 1.55, color=INK, w=1.9)
    _trapezoid(s, 4.28, y + 0.92, 1.55, 1.28, fill=TEAL, label="")
    text(s, "encoder", 4.22, y + 1.40, 1.70, 0.3, size=11.5, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "weights worth keeping", 4.15, y + 2.38, 1.90, 0.3, size=10.5,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "days of compute · done once · downloaded by everyone",
         0.85, y + h - 0.46, 5.6, 0.32, size=11, color=INK, align=PP_ALIGN.CENTER)
    chevron(s, 6.74, y + 1.36, 1.45, 0.74, "reuse", fill=ORANGE, size=10.5)
    # stage 2
    card(s, 7.72, y, 5.00, h, fill=BLUSH)
    text(s, "STAGE 2 · FINE-TUNE", 7.72, y + 0.16, 5.00, 0.3, size=13,
         color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    text(s, "many times, by you, on a laptop", 7.72, y + 0.48, 5.00, 0.3,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    _trapezoid(s, 8.30, y + 0.92, 1.40, 1.20, fill=TEAL, label="")
    text(s, "encoder", 8.24, y + 1.36, 1.55, 0.3, size=10.5, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER)
    arrow(s, 9.90, y + 1.52, 10.42, y + 1.52, color=INK, w=1.8)
    rbox(s, 10.50, y + 1.20, 1.55, 0.66, fill=WHITE, edge=ORANGE, label="your\nhead",
         size=10.5, label_color=ORANGE_D)
    scatter_cloud(s, 8.95, y + 2.62, 0.55, 0.26, 11, 55, fill=ORANGE, d=0.11)
    text(s, "10² – 10³ labelled\nmolecules", 8.05, y + 2.92, 2.0, 0.55, size=10.5,
         color=INK, align=PP_ALIGN.CENTER)
    text(s, "minutes of compute", 10.30, y + 2.92, 2.1, 0.3, size=10.5,
         color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_three_knobs(s, y=1.38, h=4.1):
    """Frozen / partial / full fine-tuning, positioned on a dataset-size axis."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    modes = [
        ("FROZEN  (linear probe)", ["frozen", "frozen", "frozen", "TRAIN"],
         "10 – 500 labels", "fast · impossible to overfit the encoder"),
        ("PARTIAL", ["frozen", "frozen", "TRAIN", "TRAIN"],
         "500 – 5,000 labels", "unfreeze the top layers only"),
        ("FULL", ["TRAIN", "TRAIN", "TRAIN", "TRAIN"],
         "> 5,000 labels", "best ceiling · needs a small learning rate"),
    ]
    for x, (hd, layers, regime, note) in zip(xs, modes):
        card(s, x, y, W, h, fill=CARD)
        text(s, hd, x, y + 0.16, W, 0.3, size=12.5, color=INK, bold=True,
             align=PP_ALIGN.CENTER)
        for k, st in enumerate(layers):
            frozen = (st == "frozen")
            rbox(s, x + 0.85, y + 0.62 + k * 0.52, W - 1.70, 0.42,
                 fill=RGBColor(0xE7, 0xE9, 0xEE) if frozen else ORANGE,
                 edge=LGRAY if frozen else ORANGE,
                 label=("frozen" if frozen else "trained"),
                 size=10.5, label_color=GRAY if frozen else WHITE, bold=not frozen)
        text(s, "input", x, y + 0.62 + 4 * 0.52 + 0.02, W, 0.26, size=9.5,
             color=GRAY, align=PP_ALIGN.CENTER)
        card(s, x + 0.30, y + h - 1.06, W - 0.60, 0.44, fill=MINT)
        text(s, regime, x + 0.30, y + h - 1.00, W - 0.60, 0.32, size=12,
             color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        text(s, note, x + 0.20, y + h - 0.54, W - 0.40, 0.42, size=10.5,
             color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_nmr_crossover(s, y=1.36, h=4.15):
    """THE key figure: pretraining pays only in the low-data regime."""
    xy_chart(s, 0.90, y + 0.28, 7.55, h - 0.80,
             [{"name": "pre-trained on DFT, then fine-tuned",
               "pts": [(97, 0.43), (386, 0.37), (676, 0.34), (966, 0.33)],
               "kind": "both", "color": TEAL, "size": 8, "width": 2.8},
              {"name": "trained from scratch on experiment",
               "pts": [(97, 1.35), (386, 0.41), (676, 0.40), (966, 0.39)],
               "kind": "both", "color": ORANGE, "size": 8, "width": 2.8,
               "dash": DASH.DASH}],
             xlabel="experimental training molecules",
             ylabel="¹H chemical-shift MAE  (ppm)",
             xlim=(0, 1050), ylim=(0, 1.5), legend=True)
    # annotate the two regimes
    px0 = 0.90 + 0.82
    pw = 7.55 - 0.98
    def PX(v):
        return px0 + v / 1050 * pw
    line(s, PX(97), y + 0.42, PX(97), y + h - 0.62, color=GRAY, w=1.2, dash=DASH.DASH)
    line(s, PX(500), y + 0.42, PX(500), y + h - 0.62, color=GRAY, w=1.2, dash=DASH.DASH)
    text(s, "3× better", PX(97) + 0.10, y + 1.08, 1.3, 0.32, size=14, color=GREEN,
         bold=True)
    text(s, "gap has closed", PX(560), y + 2.28, 2.0, 0.32, size=13, color=RED,
         bold=True)
    text(s, "~100 molecules", PX(97) - 0.85, y + h - 0.58, 1.7, 0.28, size=10,
         color=GRAY, align=PP_ALIGN.CENTER)
    # right panel: the message
    card(s, 8.75, y + 0.20, 3.97, h - 0.42, fill=MINT)
    text(s, "the whole lecture\nin one plot", 8.90, y + 0.36, 3.67, 0.66, size=15,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    pts = ["Pre-training buys you data\nyou do not have.",
           "Its value is largest when\nyou have almost nothing.",
           "It decays to zero — and for\n¹³C it went negative — once\nyou have ~1,000 labels."]
    for k, t in enumerate(pts):
        yy = y + 1.24 + k * 0.88
        text(s, "•", 9.00, yy, 0.25, 0.3, size=15, color=ORANGE, bold=True)
        text(s, t, 9.28, yy + 0.02, 3.25, 0.84, size=11.5, color=INK)
    return y + h


def fig_low_data_wins(s, y=1.4, h=4.05):
    """Three domains, same story: pretraining collapses the data requirement."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    cards = [
        ("REACTION YIELD", "98", "reactions are enough for a pre-trained\n"
         "reaction encoder to match DFT descriptors + RF",
         "Buchwald–Hartwig HTE · Yield-BERT"),
        ("INTERATOMIC POTENTIALS", "3 %", "of the data needed to specialise a "
         "foundation MLIP to liquid water,\nvs training from scratch",
         "MatterSim · 17 M structures pre-trained"),
        ("SUBLIMATION ENTHALPIES", "tens", "of training structures give sub-kJ/mol "
         "accuracy after fine-tuning MACE-MP-0",
         "Kaur et al., Faraday Discuss. 2025"),
    ]
    for x, (hd, big, body, src) in zip(xs, cards):
        card(s, x, y, W, h, fill=MINT)
        text(s, hd, x, y + 0.20, W, 0.3, size=11.5, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, big, x, y + 0.72, W, 0.95, size=54, color=ORANGE, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, body, x + 0.24, y + 1.92, W - 0.48, 1.05, size=11.5, color=INK,
             align=PP_ALIGN.CENTER)
        text(s, src, x + 0.20, y + h - 0.60, W - 0.40, 0.5, size=10, color=GRAY,
             italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_mlip_scaling(s, y=1.4, h=4.05):
    """Foundation models for atoms: the same pretrain-then-finetune story."""
    xy_chart(s, 0.90, y + 0.30, 7.30, h - 0.90,
             [{"name": "Matbench Discovery F1",
               "pts": [(5.27, 0.569), (6.20, 0.669), (7.51, 0.862), (8.12, 0.905),
                       (8.05, 0.925)],
               "kind": "markers", "color": TEAL, "size": 10}],
             xlabel="pre-training structures  (log₁₀)",
             ylabel="Matbench Discovery  F1", xlim=(5, 8.6), ylim=(0.5, 1.0))
    #                                    label                 dx     dy
    labels = [(5.27, 0.569, "M3GNet 2022",                    0.72, -0.34),
              (6.20, 0.669, "MACE-MP-0 2023",                 0.72, -0.14),
              (7.51, 0.862, "MatterSim 2024",                -0.05,  0.20),
              (8.12, 0.905, "Orb-v3 2025",                    0.30, -0.20),
              (8.05, 0.925, "eSEN 2025",                     -0.98, -0.30)]
    px0, pw = 0.90 + 0.86, 7.30 - 1.02
    py0, ph = y + 0.42, h - 0.90 - 0.58
    for lx, ly, lbl, dx, dy in labels:
        px = px0 + (lx - 5) / 3.6 * pw
        py = py0 + (1 - (ly - 0.5) / 0.5) * ph
        text(s, lbl, px - 0.95 + dx, py + dy, 1.9, 0.28,
             size=10, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "~600× more training structures in three years", 0.90, y + 0.04, 7.30,
         0.3, size=12.5, color=TEAL, bold=True, align=PP_ALIGN.LEFT)
    card(s, 8.55, y + 0.30, 4.17, h - 0.90, fill=BLUSH)
    text(s, "why chemists care", 8.70, y + 0.46, 3.87, 0.32, size=13, color=ORANGE_D,
         bold=True, align=PP_ALIGN.CENTER)
    body = ("You will almost never train one of\nthese.\n\n"
            "You download it, then fine-tune on\nthe 50–500 structures you can\n"
            "afford to compute for your own\nsystem.\n\n"
            "Exactly the recipe from the last\nslide — applied to potentials\n"
            "instead of properties.")
    text(s, body, 8.78, y + 0.92, 3.72, 2.3, size=11, color=INK)
    return y + h


# =====================================================================
# PART 6 — when it fails
# =====================================================================
def fig_benchmark_reality(s, y=1.4, h=4.05):
    """The 2025 replication: 25 models, 25 datasets, ECFP is still there."""
    card(s, 0.62, y, 6.05, h, fill=BLUSH)
    text(s, "25 pre-trained models  ×  25 datasets", 0.62, y + 0.22, 6.05, 0.34,
         size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "re-evaluated with hierarchical Bayesian testing\ninstead of "
            "leaderboard means", 0.62, y + 0.62, 6.05, 0.6, size=11, color=GRAY,
         italic=True, align=PP_ALIGN.CENTER)
    # a grid of model chips, one highlighted
    cols, rows_n = 5, 5
    gw, gh = 0.95, 0.30
    gx0 = 0.62 + (6.05 - cols * (gw + 0.10)) / 2
    gy0 = y + 1.34
    for i in range(rows_n):
        for j in range(cols):
            k = i * cols + j
            win = (k == 12)
            rbox(s, gx0 + j * (gw + 0.10), gy0 + i * (gh + 0.09), gw, gh,
                 fill=GREEN if win else RGBColor(0xE7, 0xE9, 0xEE),
                 edge=GREEN if win else RGBColor(0xE7, 0xE9, 0xEE),
                 label="CLAMP" if win else "", size=9,
                 label_color=WHITE, bold=True)
    text(s, "one model beats the ECFP baseline significantly —\nand it is itself "
            "fingerprint-based", 0.85, y + h - 0.62, 5.6, 0.58, size=11.5,
         color=RED, bold=True, align=PP_ALIGN.CENTER)
    # right: the quote
    card(s, 6.87, y, 5.85, h, fill=CARD)
    text(s, "“", 7.10, y + 0.16, 0.8, 0.7, size=46, color=LGRAY, bold=True)
    text(s, "nearly all neural models show\nnegligible or no improvement over\n"
            "the baseline ECFP molecular\nfingerprint.",
         7.25, y + 0.88, 5.15, 1.45, size=15, color=INK, italic=True)
    text(s, "Praski, Adamczyk & Czech, arXiv:2508.06199 (2025)", 7.25, y + 2.42,
         5.15, 0.3, size=10.5, color=GRAY, italic=True)
    card(s, 7.10, y + 2.86, 5.40, 1.02, fill=MINT)
    text(s, "Not a reason to give up on learned representations.\n"
            "A reason to always run the ECFP + random-forest\nbaseline first.",
         7.22, y + 2.96, 5.16, 0.86, size=12, color=TEAL, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_activity_cliff(s, y=1.4, h=4.05):
    """Where every smooth representation breaks."""
    card(s, 0.62, y, 6.90, h, fill=CARD)
    panel_title(s, "two molecules, one methyl group apart", 0.88, y + 0.14, 6.4, size=11.5)
    for k, (cx, pot, col, lbl) in enumerate([(2.30, "IC₅₀  =  8 nM", TEAL, "A"),
                                             (5.65, "IC₅₀  =  4,200 nM", RED, "B")]):
        mol_sketch(s, cx, y + 1.45, sc=1.75, tails=3, seed=4, color=INK,
                   hub=ORANGE if k else None)
        text(s, pot, cx - 1.30, y + 2.52, 2.6, 0.34, size=14, color=col, bold=True,
             align=PP_ALIGN.CENTER)
    text(s, "Tanimoto similarity  0.95", 0.88, y + 2.95, 6.4, 0.32, size=12.5,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "500× potency difference", 0.88, y + 3.26, 6.4, 0.32, size=12.5,
         color=RED, bold=True, align=PP_ALIGN.CENTER)
    # right
    card(s, 7.72, y, 5.00, h, fill=BLUSH)
    text(s, "what breaks", 7.72, y + 0.18, 5.00, 0.32, size=13.5, color=RED,
         bold=True, align=PP_ALIGN.CENTER)
    pts = ["Smooth embeddings put A and B\nin the same place.",
           "Errors on cliff compounds reach\n1.4 log units — a 25× error in\npotency.",
           "Across 30 targets, SVM on ECFP\nwas the best method overall.",
           "Transformers pre-trained on 10 M\nSMILES did not beat LSTMs\n"
           "pre-trained on 36 k."]
    for k, t in enumerate(pts):
        text(s, "•", 7.95, y + 0.66 + k * 0.78, 0.25, 0.3, size=15, color=RED, bold=True)
        text(s, t, 8.22, y + 0.68 + k * 0.78, 4.35, 0.74, size=11, color=INK)
    return y + h


def fig_crossover(s, y=1.38, h=4.1):
    """When is a learned representation worth it? Schematic, honestly labelled."""
    xy_chart(s, 0.90, y + 0.30, 7.85, h - 0.85,
             [{"name": "ECFP + random forest",
               "pts": [(2.0, 0.62), (2.5, 0.68), (3.0, 0.73), (3.5, 0.765),
                       (4.0, 0.785), (4.5, 0.795), (5.0, 0.80)],
               "kind": "line", "color": ORANGE, "width": 3.0, "smooth": True},
              {"name": "deep model, from scratch",
               "pts": [(2.0, 0.30), (2.5, 0.42), (3.0, 0.58), (3.5, 0.72),
                       (4.0, 0.83), (4.5, 0.89), (5.0, 0.92)],
               "kind": "line", "color": TEAL, "width": 3.0, "smooth": True},
              {"name": "pre-trained encoder, fine-tuned",
               "pts": [(2.0, 0.58), (2.5, 0.68), (3.0, 0.77), (3.5, 0.84),
                       (4.0, 0.89), (4.5, 0.92), (5.0, 0.94)],
               "kind": "line", "color": PURPLE, "width": 3.0, "smooth": True}],
             xlabel="labelled training molecules  (log₁₀)",
             ylabel="predictive performance", xlim=(2, 5), ylim=(0.2, 1.0),
             legend=True, tick_labels=True)
    px0, pw = 0.90 + 0.88, 7.85 - 1.05
    def PX(v):
        return px0 + (v - 2) / 3 * pw
    for v, lbl, col in [(3.0, "≈ 1,000", GRAY), (4.0, "≈ 10,000", GRAY)]:
        line(s, PX(v), y + 0.44, PX(v), y + h - 0.62, color=col, w=1.1,
             dash=DASH.DASH)
    for lbl, cx, col, tint in [("fingerprints win", PX(2.42), ORANGE_D, BLUSH),
                               ("pre-training earns its keep", PX(3.50), PURPLE,
                                RGBColor(0xF2, 0xEC, 0xF4)),
                               ("everything converges", PX(4.55), TEAL, MINT)]:
        card(s, cx - 1.05, y + h - 0.62, 2.10, 0.44, fill=tint)
        text(s, lbl, cx - 1.05, y + h - 0.56, 2.10, 0.34, size=11, color=col,
             bold=True, align=PP_ALIGN.CENTER)
    card(s, 9.05, y + 0.24, 3.67, h - 0.62, fill=BLUSH)
    text(s, "schematic", 9.20, y + 0.36, 3.37, 0.3, size=11.5, color=RED,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "The shapes are real;\nthe exact crossover is not.\n\n"
            "Published estimates put it\nanywhere between 10³ and 10⁴\n"
            "labelled molecules, and it moves\nwith label noise, split type and\n"
            "how narrow your chemistry is.\n\n"
            "Measure it on your own data —\nthat is what the tutorial does.",
         9.22, y + 0.76, 3.33, 2.8, size=11, color=INK)
    return y + h


def fig_decision(s, y=1.34, h=4.15):
    """A short decision path the audience can actually use on Monday."""
    q = [("How many labelled examples do you have?", None),
         ("< 100", "Do not train a deep model.\nECFP + RF, or a Gaussian process.\n"
                   "Consider frozen pre-trained embeddings\nas extra features."),
         ("100 – 5,000", "This is where pre-training pays.\n"
                         "Frozen embeddings first, then partial\nfine-tuning. "
                         "Always report the ECFP baseline."),
         ("> 10,000", "Train from scratch on your own data.\n"
                      "Pre-training becomes a convenience,\nnot a necessity.")]
    text(s, q[0][0], 0.62, y, SW - 1.24, 0.4, size=16, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    accents = [ORANGE, TEAL, PURPLE]
    tints = [BLUSH, MINT, RGBColor(0xF2, 0xEC, 0xF4)]
    for x, (lbl, body), acc, tint in zip(xs, q[1:], accents, tints):
        arrow(s, SW / 2, y + 0.46, x + W / 2, y + 0.86, color=GRAY, w=1.4)
        card(s, x, y + 0.90, W, h - 1.55, fill=tint)
        text(s, lbl, x, y + 1.06, W, 0.36, size=17, color=acc, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, body, x + 0.30, y + 1.60, W - 0.60, 1.6, size=12, color=INK,
             align=PP_ALIGN.LEFT, line_spacing=1.25)
    card(s, 0.62, y + h - 0.58, SW - 1.24, 0.58, fill=PEACH)
    text(s, "In all three columns: run the baseline first, and split by scaffold, "
            "not at random.", 0.62, y + h - 0.50, SW - 1.24, 0.42, size=13.5,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    return y + h
