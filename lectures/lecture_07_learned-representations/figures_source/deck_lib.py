#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Slide helpers styled to match lecture_02_chemical-representations:
  · 13.333 x 7.5 in
  · peach title bar (0 -> 1.023 in), Avenir Book 28pt centered, black
  · peach footer bar (7.023 -> 7.523 in), right-aligned "AI4ChemicalSciences BootCamp 2026"
  · Caltech wordmark bottom-left
"""
import math, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.xmlchemy import OxmlElement

# ---------------- palette (matched to the template theme) ----------------
PEACH    = RGBColor(0xFB, 0xE4, 0xD6)   # title / footer bar
ORANGE   = RGBColor(0xE9, 0x71, 0x32)   # theme accent2 (Caltech orange)
ORANGE_D = RGBColor(0xB4, 0x55, 0x1F)
TEAL     = RGBColor(0x15, 0x60, 0x82)   # theme accent1
TEAL_L   = RGBColor(0x4E, 0x9A, 0xB8)
GREEN    = RGBColor(0x19, 0x6B, 0x24)
PURPLE   = RGBColor(0xA0, 0x2B, 0x93)
INK      = RGBColor(0x20, 0x24, 0x2E)
GRAY     = RGBColor(0x6E, 0x74, 0x80)
LGRAY    = RGBColor(0xD9, 0xDC, 0xE3)
CARD     = RGBColor(0xF4, 0xF6, 0xFA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Avenir Book"          # template font
FONT_B = "Avenir Book"

SW, SH = 13.3333, 7.5
BAR_H     = 1.023
FOOT_Y    = 7.023
FOOT_H    = 0.5
BODY_TOP  = BAR_H + 0.18
BODY_BOT  = FOOT_Y - 0.12
LOGO = "caltech_logo.png"
FOOTER_TEXT = "AI4ChemicalSciences BootCamp 2026"

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
_count = [0]

# ---------------- primitives ----------------
def _margins(tf):
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True

def text(slide, runs, x, y, w, h, size=14, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
         space_after=0, font=FONT):
    """runs: str, or list of (txt, bold, color) for mixed formatting in one paragraph."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    _margins(tf)
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for t, b, c in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b; r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = c if c is not None else color
    return box

def _parse_bold(s):
    """'**lead** rest' -> [(lead, True, None), (rest, False, None)].
    Single *...* is treated as emphasis but rendered plain (no literal asterisks)."""
    s = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", s)
    out, pos = [], 0
    for m in re.finditer(r"\*\*(.+?)\*\*", s):
        if m.start() > pos:
            out.append((s[pos:m.start()], False, None))
        out.append((m.group(1), True, None))
        pos = m.end()
    if pos < len(s):
        out.append((s[pos:], False, None))
    return out or [(s, False, None)]

def _lines(s, w_in, size):
    """Conservative wrap estimate calibrated against the rendered deck (Avenir/DejaVu)."""
    cpl = max(10, int(w_in * (16.5 / size) * 7.2))
    return max(1, math.ceil(len(s) / cpl))

def new_slide(title=None, footer=True, title_size=28):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    _count[0] += 1
    # title bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(BAR_H))
    bar.fill.solid(); bar.fill.fore_color.rgb = PEACH
    bar.line.fill.background(); bar.shadow.inherit = False
    if title:
        tf = bar.text_frame
        _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sz = title_size if len(title) < 56 else (23 if len(title) < 76 else 20)
        r = p.add_run(); r.text = title
        r.font.size = Pt(sz); r.font.name = FONT; r.font.color.rgb = INK
    # footer bar
    if footer:
        fb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(FOOT_Y),
                                Inches(SW), Inches(FOOT_H))
        fb.fill.solid(); fb.fill.fore_color.rgb = PEACH
        fb.line.fill.background(); fb.shadow.inherit = False
        tf = fb.text_frame
        _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_right = Inches(0.28)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = FOOTER_TEXT
        r.font.size = Pt(14); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = INK
        try:
            s.shapes.add_picture(LOGO, Inches(0.05), Inches(7.07),
                                 Inches(1.70), Inches(0.41))
        except Exception:
            pass
    return s

def bullets(slide, items, x, y, w, size=16, color=INK, gap=0.20, marker="dot",
            marker_color=ORANGE, start_num=1, lh=0.27):
    """Literal-bullet style matching the template ('• text'), or numbered circles."""
    cur = y
    for i, raw in enumerate(items):
        plain = raw.replace("**", "")
        n = _lines(plain, w - 0.45, size)
        h = max(0.32, n * lh + 0.10)
        if marker == "dot":
            text(slide, "•", x, cur - 0.015, 0.3, 0.4, size=size + 2,
                 color=marker_color, bold=True)
            tx = x + 0.28
        elif marker == "num":
            d = 0.32
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cur + 0.02),
                                       Inches(d), Inches(d))
            c.fill.solid(); c.fill.fore_color.rgb = marker_color
            c.line.fill.background(); c.shadow.inherit = False
            tf = c.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(start_num + i)
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
            r.font.name = FONT
            tx = x + 0.45
        else:
            tx = x
        runs = [(t, b, color if c is None else c) for t, b, c in _parse_bold(raw)]
        text(slide, runs, tx, cur, w - (tx - x), h + 0.12, size=size, color=color)
        cur += h + gap
    return cur

try:
    from PIL import Image as _PILImage
except Exception:
    _PILImage = None

def img_aspect(path):
    """height / width of the image file."""
    if _PILImage is None:
        return 0.5
    with _PILImage.open(path) as im:
        return im.size[1] / im.size[0]

def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def place_image(slide, path, y, w, x=None, max_bottom=None):
    """Place an image by width, measuring its true height. Shrinks to respect
    max_bottom. Returns the actual bottom edge in inches."""
    asp = img_aspect(path)
    h = w * asp
    if max_bottom is not None and y + h > max_bottom:
        h = max_bottom - y
        w = h / asp
    if x is None:
        x = (SW - w) / 2
    picture(slide, path, x, y, w=w)
    return y + h

def fig_slide(title, img, caption=None, notes=None, img_w=11.4, img_top=1.35,
              note_size=14, title_size=28, note_gap=0.12, note_lh=0.25):
    """Big centred figure with optional bullets underneath.
    The figure is shrunk if necessary so the notes always fit above the footer."""
    s = new_slide(title, title_size=title_size)
    # reserve room for the notes block
    reserve = 0.30
    if notes:
        for t in notes:
            reserve += max(0.32, _lines(t.replace("**", ""), SW - 1.95, note_size)
                           * note_lh + 0.10) + note_gap
    if caption:
        reserve += 0.38
    bottom = place_image(s, img, img_top, img_w, max_bottom=BODY_BOT - reserve)
    if caption:
        text(s, caption, 0.8, bottom + 0.06, SW - 1.6, 0.35, size=12,
             color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        bottom += 0.38
    if notes:
        bullets(s, notes, 0.75, bottom + 0.14, SW - 1.5, size=note_size,
                gap=note_gap, lh=note_lh)
    return s

def card(slide, x, y, w, h, fill=CARD, line=None):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    c.adjustments[0] = 0.06
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if line is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = line; c.line.width = Pt(1.25)
    c.shadow.inherit = False
    return c

def key_box(slide, msg, y=None, fill=PEACH, size=15, w=None):
    """Take-home strip across the bottom of the body area.
    Auto-sizes to the text and is clamped so it can never cross into the footer."""
    w = w or (SW - 1.4)
    plain = msg.replace("**", "")
    n = _lines(plain, w - 0.5, size)
    h = max(0.62, 0.30 + n * 0.30)
    if y is None:
        y = BODY_BOT - h
    # never overlap the footer: push the box up if it would cross BODY_BOT
    if y + h > BODY_BOT:
        y = BODY_BOT - h
    card(slide, (SW - w) / 2, y, w, h, fill=fill)
    runs = [(t, b, TEAL if b else INK) for t, b, c in _parse_bold(msg)]
    text(slide, runs, (SW - w) / 2 + 0.25, y + 0.06, w - 0.5, h - 0.12, size=size,
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    return y

def section_slide(num, title, items):
    s = new_slide(None)
    text(s, f"{num:02d}", 0.9, 1.75, 2.2, 1.5, size=76, color=PEACH, bold=True)
    text(s, title, 0.95, 3.05, 11.5, 1.0, size=32, color=INK, line_spacing=1.05)
    y = 4.25
    for it in items:
        text(s, "•", 1.0, y, 0.3, 0.35, size=17, color=ORANGE, bold=True)
        text(s, it, 1.30, y, 10.5, 0.4, size=15, color=GRAY)
        y += 0.44
    return s

def table(slide, headers, rows, x=0.75, y=1.5, w=None, col_w=None, size=13,
          head_size=13.5, row_h=None):
    w = w or (SW - 2 * x)
    n_r, n_c = len(rows) + 1, len(headers)
    h = row_h or min(4.6, 0.52 + 0.52 * len(rows))
    gt = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Inches(cw)
    for j, htxt in enumerate(headers):
        cell = gt.cell(0, j); cell.text = htxt
        cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(head_size); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gt.cell(i + 1, j); cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i % 2 == 0 else WHITE
            np_ = len(cell.text_frame.paragraphs)
            for pi, p in enumerate(cell.text_frame.paragraphs):
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                p.line_spacing = 1.05
                last = (pi == np_ - 1)
                for r in p.runs:
                    r.font.name = FONT
                    if np_ > 1:
                        r.font.size = Pt(size + 1 if last else size - 2)
                        r.font.bold = (j == 0) or last
                        r.font.color.rgb = TEAL if j == 0 else (INK if last else GRAY)
                    else:
                        r.font.size = Pt(size)
                        r.font.bold = (j == 0)
                        r.font.color.rgb = TEAL if j == 0 else INK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
    return gt

# ---------------- vector diagram primitives (for perceptron build-up) ----------------
def line(slide, x1, y1, x2, y2, color=GRAY, w=1.25, arrow=False, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    if dash: c.line.dash_style = dash
    c.shadow.inherit = False
    if arrow:
        ln = c.line._get_or_add_ln()
        t = OxmlElement('a:tailEnd')
        t.set('type', 'triangle'); t.set('w', 'med'); t.set('len', 'med')
        ln.append(t)
    return c

def node(slide, cx, cy, d, fill=WHITE, edge=INK, label="", size=12,
         label_color=INK, bold=True, edge_w=1.5):
    n = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2),
                               Inches(d), Inches(d))
    n.fill.solid(); n.fill.fore_color.rgb = fill
    n.line.color.rgb = edge; n.line.width = Pt(edge_w)
    n.shadow.inherit = False
    if label:
        tf = n.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = label_color; r.font.name = FONT
    return n

def rbox(slide, x, y, w, h, fill=WHITE, edge=INK, label="", size=12,
         label_color=INK, bold=True):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    b.adjustments[0] = 0.16
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = edge; b.line.width = Pt(1.4)
    b.shadow.inherit = False
    if label:
        tf = b.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, ln_ in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = ln_
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = label_color; r.font.name = FONT
    return b

def centered(slide, txt, y, size=15, color=INK, bold=False, italic=False, w=None):
    w = w or (SW - 1.6)
    return text(slide, txt, (SW - w) / 2, y, w, 0.42, size=size, color=color,
                bold=bold, italic=italic, align=PP_ALIGN.CENTER)

def save(path="deck.pptx"):
    prs.save(path)
    print("slides:", _count[0], "->", path)
