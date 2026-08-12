#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Figures for lecture_07 "Learned Representations".

Everything here is drawn with NATIVE PowerPoint objects:
  · data plots  -> real charts   (right-click > Edit Data works)
  · schematics  -> real autoshapes + connectors (click, drag, recolour)
  · tables      -> real PowerPoint tables
No raster images anywhere, so every figure stays editable in Keynote/PowerPoint.
"""
import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH

from deck_lib import (PEACH, ORANGE, ORANGE_D, TEAL, TEAL_L, GREEN, PURPLE,
                      INK, GRAY, LGRAY, CARD, WHITE, FONT, SW, SH,
                      BODY_TOP, BODY_BOT,
                      text, node, rbox, line, card, centered, _margins)
from chartlib import xy_chart, col_chart, panel_title, caption

# extra tints used only in this deck
BLUSH  = RGBColor(0xFD, 0xF1, 0xE8)   # very light peach card
MINT   = RGBColor(0xE8, 0xF4, 0xF0)   # very light teal card
SAND   = RGBColor(0xF7, 0xF3, 0xEA)
RED    = RGBColor(0xC0, 0x39, 0x2B)
OLIVE  = RGBColor(0x7A, 0x8B, 0x2E)


# =====================================================================
# small shared helpers
# =====================================================================
def notes(slide, txt):
    """Attach speaker notes (what to say) to a slide."""
    slide.notes_slide.notes_text_frame.text = txt.strip()
    return slide


def arrow(slide, x1, y1, x2, y2, color=GRAY, w=1.8, dash=None):
    return line(slide, x1, y1, x2, y2, color=color, w=w, arrow=True, dash=dash)


def chevron(slide, x, y, w, h, label, fill=TEAL, label_color=WHITE, size=12):
    sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = (i == 0)
        r.font.color.rgb = label_color; r.font.name = FONT
    return sh


def blob(slide, cx, cy, rx, ry, fill=MINT, edge=None, alpha_line=True):
    """Soft ellipse used for latent-space clouds and cluster hulls."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - rx), Inches(cy - ry),
                                Inches(2 * rx), Inches(2 * ry))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if edge is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = edge; sh.line.width = Pt(1.2)
        if alpha_line:
            sh.line.dash_style = DASH.DASH
    sh.shadow.inherit = False
    return sh


def dot(slide, cx, cy, d=0.10, fill=TEAL, edge=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2),
                                Inches(cy - d / 2), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if edge is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = edge; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _rng(seed):
    """Tiny deterministic LCG so figures are reproducible without numpy."""
    s = [seed]
    def r():
        s[0] = (1103515245 * s[0] + 12345) % (2 ** 31)
        return s[0] / (2 ** 31)
    return r


def scatter_cloud(slide, cx, cy, rx, ry, n, seed, fill=TEAL, d=0.09,
                  skip=None):
    """n dots scattered in an ellipse. `skip(x, y) -> True` punches holes."""
    r = _rng(seed)
    out = []
    tries = 0
    while len(out) < n and tries < n * 40:
        tries += 1
        a, b = r() * 2 - 1, r() * 2 - 1
        if a * a + b * b > 1:
            continue
        x, y = cx + a * rx, cy + b * ry
        if skip and skip(a, b):
            continue
        dot(slide, x, y, d=d, fill=fill)
        out.append((x, y))
    return out


def mol_sketch(slide, cx, cy, sc=1.0, ring=True, tails=2, seed=3,
               color=INK, hub=None, label=None):
    """A generic little molecule cartoon: one hexagon plus a couple of tails.
    Purely decorative but drawn as real shapes, so it can be restyled."""
    pts = []
    for k in range(6):
        th = math.radians(60 * k - 90)
        pts.append((cx + 0.30 * sc * math.cos(th), cy + 0.30 * sc * math.sin(th)))
    if ring:
        for k in range(6):
            a, b = pts[k], pts[(k + 1) % 6]
            line(slide, a[0], a[1], b[0], b[1], color=color, w=1.6 * sc)
    r = _rng(seed)
    tip = []
    for t in range(tails):
        k = int(r() * 6) % 6
        a = pts[k]
        th = math.atan2(a[1] - cy, a[0] - cx)
        bx, by = a[0] + 0.26 * sc * math.cos(th), a[1] + 0.26 * sc * math.sin(th)
        line(slide, a[0], a[1], bx, by, color=color, w=1.6 * sc)
        tip.append((bx, by))
    if hub:
        for p in tip[:1] or [pts[0]]:
            dot(slide, p[0], p[1], d=0.13 * sc, fill=hub)
    if label:
        text(slide, label, cx - 0.7 * sc, cy + 0.42 * sc, 1.4 * sc, 0.26,
             size=9.5 * sc, color=GRAY, align=PP_ALIGN.CENTER)
    return pts


def stack(slide, x, y, w, layers, lh=0.30, gap=0.055, size=10.5):
    """A vertical stack of labelled layer bars (used for encoder diagrams)."""
    cur = y
    boxes = []
    for lbl, fill, fg in layers:
        b = rbox(slide, x, cur, w, lh, fill=fill, edge=fill, label=lbl,
                 size=size, label_color=fg, bold=False)
        boxes.append(b)
        cur += lh + gap
    return cur - gap


def legend(slide, items, x, y, size=11, gap=0.30, horizontal=True, sw=0.16):
    """items: [(label, color), ...]"""
    cx, cy = x, y
    for lbl, col in items:
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy + 0.045),
                                    Inches(sw), Inches(sw))
        sh.fill.solid(); sh.fill.fore_color.rgb = col
        sh.line.fill.background(); sh.shadow.inherit = False
        w = 0.10 + 0.075 * size * 0.115 * len(lbl)
        text(slide, lbl, cx + sw + 0.09, cy, w, 0.26, size=size, color=INK)
        if horizontal:
            cx += sw + 0.09 + w + gap
        else:
            cy += 0.30


def log_value_axis(ch, base=10):
    """python-pptx has no log-scale property; set c:logBase on the value axis.
    Needed whenever one series is 100x another (e.g. TMAP vs UMAP runtimes)."""
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    ax = ch.value_axis._element
    scaling = ax.find(qn('c:scaling'))
    if scaling is None:
        return ch
    for tag in ('c:min', 'c:max'):                 # log + explicit min/max clash
        for e in scaling.findall(qn(tag)):
            scaling.remove(e)
    for e in scaling.findall(qn('c:logBase')):
        scaling.remove(e)
    lb = OxmlElement('c:logBase')
    lb.set('val', str(base))
    scaling.insert(0, lb)                          # c:logBase precedes c:orientation
    return ch


def ref(slide, txt, y=None):
    """Small grey citation line pinned just above the footer."""
    y = y if y is not None else BODY_BOT - 0.02
    return text(slide, txt, 0.62, y, SW - 1.24, 0.26, size=10.5, color=GRAY,
                italic=True)


# =====================================================================
# PART 1 — why learn a representation
# =====================================================================
def fig_two_pipelines(s, y=1.42):
    """Hand-crafted features vs learned features, as two parallel pipelines."""
    rows = [
        (y + 0.00, "YOU choose the features", ORANGE, BLUSH,
         ["molecule", "rules you\nwrote", "fixed vector\n(ECFP, descriptors)", "small model", "property"]),
        (y + 2.30, "THE MODEL chooses the features", TEAL, MINT,
         ["molecule", "encoder\n(learned)", "embedding\n(learned vector)", "head", "property"]),
    ]
    bw, bh, gapx = 2.05, 1.00, 0.42
    x0 = 1.62
    for ry, title, accent, tint, labels in rows:
        text(s, title, 0.62, ry - 0.30, 5.6, 0.28, size=13, color=accent, bold=True)
        card(s, 0.55, ry - 0.06, SW - 1.10, bh + 0.12, fill=tint)
        for i, lbl in enumerate(labels):
            bx = x0 + i * (bw + gapx)
            fill = WHITE
            edge = accent if i in (1, 2) else LGRAY
            rbox(s, bx, ry, bw, bh, fill=fill, edge=edge, label=lbl,
                 size=11.5, label_color=INK, bold=(i in (1, 2)))
            if i < len(labels) - 1:
                arrow(s, bx + bw + 0.05, ry + bh / 2,
                      bx + bw + gapx - 0.05, ry + bh / 2, color=accent, w=1.8)
        # brace under the learned part
    text(s, "fixed before you see the task", x0 + 1 * (bw + gapx), y + 1.06,
         bw * 2 + gapx, 0.26, size=10.5, color=ORANGE_D, italic=True,
         align=PP_ALIGN.CENTER)
    text(s, "fitted from data — this is the learned representation",
         x0 + 1 * (bw + gapx), y + 3.36, bw * 2 + gapx, 0.26, size=10.5,
         color=TEAL, italic=True, align=PP_ALIGN.CENTER)
    return y + 3.7


def fig_label_gap(s, y=1.4, h=4.05):
    """Log-scale column chart: how many molecules exist vs how many are labelled.
    Counts live in the category labels so nothing has to be positioned by hand."""
    cats = ["GDB-17\nenumerated\n1.7 × 10¹¹", "Enamine REAL\norderable\n9.6 × 10⁹",
            "PubChem\ncatalogued\n1.2 × 10⁸", "ChEMBL\nbioactivity\n2.4 × 10⁶",
            "MoleculeNet\ntypical task\n~2 × 10³", "your HTE\nplate\n~4 × 10²"]
    vals = [11.22, 9.98, 8.08, 6.38, 3.30, 2.60]          # log10 counts
    colors = [[LGRAY, LGRAY, TEAL_L, TEAL, ORANGE, ORANGE_D]]
    col_chart(s, 1.05, y, SW - 2.10, h - 0.30, cats, [("log10 molecules", vals)],
              ylabel="molecules  (log₁₀)", ylim=(0, 12), colors=colors,
              gap=45, tick_size=10.5)
    card(s, 8.15, y + 0.18, 3.95, 0.92, fill=BLUSH)
    text(s, "8 orders of magnitude", 8.15, y + 0.30, 3.95, 0.34, size=14,
         color=RED, bold=True, align=PP_ALIGN.CENTER)
    text(s, "between what we can draw\nand what we have measured",
         8.15, y + 0.64, 3.95, 0.44, size=11, color=INK, align=PP_ALIGN.CENTER)
    return y + h


def fig_supervision(s, y=1.4, h=4.05):
    """Supervised / unsupervised / self-supervised, side by side."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    titles = ["SUPERVISED", "UNSUPERVISED", "SELF-SUPERVISED"]
    subs = ["learn x → y", "learn the shape of x", "learn x → part of x"]
    accents = [ORANGE, TEAL, PURPLE]
    tints = [BLUSH, MINT, RGBColor(0xF2, 0xEC, 0xF4)]
    for i, (x, t, sub, acc, tint) in enumerate(zip(xs, titles, subs, accents, tints)):
        card(s, x, y, W, h, fill=tint)
        text(s, t, x, y + 0.16, W, 0.3, size=13, color=acc, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, sub, x, y + 0.48, W, 0.3, size=11.5, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)
        cy = y + 1.62
        if i == 0:                                    # labelled points
            r = _rng(11)
            for k in range(18):
                px = x + 0.55 + r() * (W - 1.1)
                py = cy - 0.75 + r() * 1.5
                dot(s, px, py, d=0.15, fill=ORANGE if r() > 0.5 else TEAL)
            text(s, "every point carries a label", x, y + h - 1.02, W, 0.3,
                 size=11, color=INK, align=PP_ALIGN.CENTER)
            text(s, "labels are expensive", x, y + h - 0.70, W, 0.3, size=11,
                 color=RED, bold=True, align=PP_ALIGN.CENTER)
        elif i == 1:                                  # bare cloud + clusters
            blob(s, x + 1.25, cy - 0.28, 0.62, 0.46, fill=WHITE, edge=TEAL)
            blob(s, x + 2.60, cy + 0.30, 0.70, 0.42, fill=WHITE, edge=TEAL)
            scatter_cloud(s, x + 1.25, cy - 0.28, 0.50, 0.34, 12, 21, fill=GRAY, d=0.13)
            scatter_cloud(s, x + 2.60, cy + 0.30, 0.56, 0.30, 12, 22, fill=GRAY, d=0.13)
            text(s, "no labels at all — only structure", x, y + h - 1.02, W, 0.3,
                 size=11, color=INK, align=PP_ALIGN.CENTER)
            text(s, "unlimited data", x, y + h - 0.70, W, 0.3, size=11,
                 color=GREEN, bold=True, align=PP_ALIGN.CENTER)
        else:                                         # masked token
            toks = ["C", "C", "?", "C", "=", "O"]
            tw = 0.44
            bx = x + (W - len(toks) * (tw + 0.06)) / 2
            for k, tk in enumerate(toks):
                hid = (tk == "?")
                rbox(s, bx + k * (tw + 0.06), cy - 0.72, tw, 0.44,
                     fill=PURPLE if hid else WHITE, edge=PURPLE if hid else LGRAY,
                     label=tk, size=12, label_color=WHITE if hid else INK)
            arrow(s, x + W / 2, cy - 0.20, x + W / 2, cy + 0.22, color=PURPLE)
            rbox(s, x + W / 2 - 0.85, cy + 0.26, 1.70, 0.44, fill=WHITE,
                 edge=PURPLE, label="predict “O”", size=11.5, label_color=PURPLE)
            text(s, "the data labels itself", x, y + h - 1.02, W, 0.3, size=11,
                 color=INK, align=PP_ALIGN.CENTER)
            text(s, "supervised signal, no annotation", x, y + h - 0.70, W, 0.3,
                 size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    return y + h


# =====================================================================
# PART 2 — unsupervised learning
# =====================================================================
def fig_three_jobs(s, y=1.42, h=4.0):
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    for i, x in enumerate(xs):
        card(s, x, y, W, h, fill=CARD)
    # 1 clustering
    x = xs[0]
    text(s, "1 · CLUSTERING", x, y + 0.16, W, 0.3, size=12.5, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "which molecules belong together?", x, y + 0.46, W, 0.28, size=10.5,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    for (bx, by, br, col) in [(x + 1.15, y + 1.55, 0.52, ORANGE),
                              (x + 2.62, y + 1.32, 0.46, TEAL),
                              (x + 1.95, y + 2.55, 0.48, PURPLE)]:
        blob(s, bx, by, br, br * 0.80, fill=WHITE, edge=col)
        scatter_cloud(s, bx, by, br * 0.72, br * 0.58, 8, int(bx * 97), fill=col, d=0.12)
    text(s, "groups, no coordinates", x, y + h - 0.44, W, 0.3, size=11,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    # 2 dimensionality reduction
    x = xs[1]
    text(s, "2 · DIMENSIONALITY REDUCTION", x, y + 0.16, W, 0.3, size=12.5,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "can I see 2048 dimensions?", x, y + 0.46, W, 0.28, size=10.5,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    rbox(s, x + 0.35, y + 1.10, 0.62, 1.60, fill=WHITE, edge=LGRAY, label="")
    for k in range(10):
        line(s, x + 0.42, y + 1.22 + k * 0.155, x + 0.90, y + 1.22 + k * 0.155,
             color=TEAL_L if k % 3 else ORANGE, w=3.2)
    text(s, "2048-D", x + 0.16, y + 2.76, 1.0, 0.26, size=10, color=GRAY,
         align=PP_ALIGN.CENTER)
    arrow(s, x + 1.10, y + 1.90, x + 1.62, y + 1.90, color=INK, w=2.0)
    rbox(s, x + 1.72, y + 1.10, 1.85, 1.60, fill=WHITE, edge=LGRAY, label="")
    scatter_cloud(s, x + 2.64, y + 1.90, 0.72, 0.58, 26, 44, fill=TEAL, d=0.11)
    text(s, "2-D", x + 2.15, y + 2.76, 1.0, 0.26, size=10, color=GRAY,
         align=PP_ALIGN.CENTER)
    text(s, "coordinates you can plot", x, y + h - 0.44, W, 0.3, size=11,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    # 3 density / generation
    x = xs[2]
    text(s, "3 · DENSITY & GENERATION", x, y + 0.16, W, 0.3, size=12.5,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(s, "what does a plausible molecule look like?", x, y + 0.46, W, 0.28,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    for k, (rx, ry, col) in enumerate([(1.30, 1.02, RGBColor(0xE6, 0xF0, 0xF5)),
                                       (0.95, 0.74, RGBColor(0xC9, 0xE0, 0xEB)),
                                       (0.60, 0.46, TEAL_L)]):
        blob(s, x + W / 2, y + 1.95, rx, ry, fill=col)
    dot(s, x + W / 2 - 0.10, y + 1.88, d=0.17, fill=ORANGE)
    text(s, "★", x + W / 2 + 1.02, y + 2.52, 0.4, 0.3, size=15, color=RED, bold=True)
    text(s, "sample new points from the model", x, y + h - 0.44, W, 0.3, size=11,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_pca(s, y=1.4, h=4.05):
    """A cloud with its two principal axes, plus the projection onto PC1."""
    cx, cy = 4.15, y + 1.72
    card(s, 0.70, y, 6.30, h, fill=CARD)
    r = _rng(7)
    pts = []
    for k in range(46):
        u = (r() * 2 - 1) + (r() * 2 - 1)
        v = (r() * 2 - 1) * 0.42
        # rotate 30 deg
        a, b = 0.866 * u - 0.5 * v, 0.5 * u + 0.866 * v
        px, py = cx + a * 1.20, cy - b * 0.92
        pts.append((px, py, a))
        dot(s, px, py, d=0.11, fill=TEAL)
    # PC axes
    L1 = 1.95
    line(s, cx - 0.866 * L1, cy + 0.5 * L1, cx + 0.866 * L1, cy - 0.5 * L1,
         color=ORANGE, w=2.6)
    line(s, cx + 0.5 * 0.85, cy + 0.866 * 0.85, cx - 0.5 * 0.85, cy - 0.866 * 0.85,
         color=PURPLE, w=2.2, dash=DASH.DASH)
    text(s, "PC1 — most variance", cx + 0.62, cy - 1.52, 2.4, 0.28, size=11.5,
         color=ORANGE, bold=True)
    text(s, "PC2", cx - 1.10, cy - 1.20, 0.9, 0.28, size=11.5, color=PURPLE, bold=True)
    text(s, "a rotation, nothing more", 0.70, y + h - 0.42, 6.30, 0.3, size=11.5,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    # scree
    col_chart(s, 7.35, y + 0.30, 5.30, h - 0.95,
              ["PC1", "PC2", "PC3", "PC4", "PC5", "PC6", "PC7", "PC8"],
              [("variance explained", [41, 19, 11, 7, 5, 4, 3, 2])],
              ylabel="variance explained  (%)", ylim=(0, 45),
              colors=[[ORANGE, PURPLE, TEAL, TEAL, TEAL_L, TEAL_L, LGRAY, LGRAY]],
              gap=55, tick_size=10)
    text(s, "2 components ≈ 60% of the variance — the rest is not noise, it is chemistry",
         7.35, y + h - 0.42, 5.30, 0.3, size=10.5, color=RED, italic=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_neighbour_embedding(s, y=1.42, h=3.95):
    """t-SNE / UMAP: preserve who-is-near-whom, not how-far."""
    # left: high-D neighbourhoods
    card(s, 0.62, y, 5.55, h, fill=CARD)
    panel_title(s, "high-dimensional fingerprint space", 0.85, y + 0.14, 5.2, size=11.5)
    anchors = [(2.05, y + 1.25, ORANGE), (4.35, y + 1.08, TEAL), (3.10, y + 2.35, PURPLE)]
    for (ax, ay, col) in anchors:
        blob(s, ax, ay, 0.62, 0.46, fill=WHITE, edge=col)
        scatter_cloud(s, ax, ay, 0.46, 0.32, 9, int(ax * 131), fill=col, d=0.12)
    text(s, "distances here are mostly meaningless\n(everything is far from everything)",
         0.85, y + h - 0.80, 5.1, 0.6, size=11, color=RED, italic=True,
         align=PP_ALIGN.CENTER)
    chevron(s, 6.24, y + 1.28, 1.40, 0.90, "t-SNE\nUMAP", fill=INK, size=10)
    # right: 2-D map
    card(s, 7.22, y, 5.50, h, fill=MINT)
    panel_title(s, "2-D neighbour-preserving map", 7.45, y + 0.14, 5.1, size=11.5)
    for (ax, ay, col) in [(8.75, y + 1.22, ORANGE), (11.35, y + 1.45, TEAL),
                          (9.85, y + 2.35, PURPLE)]:
        blob(s, ax, ay, 0.58, 0.44, fill=WHITE, edge=col)
        scatter_cloud(s, ax, ay, 0.42, 0.30, 9, int(ax * 77), fill=col, d=0.12)
    text(s, "neighbours stay neighbours ✓", 7.45, y + h - 0.80, 5.1, 0.3,
         size=11.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    text(s, "gap widths and cluster sizes are not meaningful ✗", 7.45,
         y + h - 0.48, 5.1, 0.3, size=11.5, color=RED, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_map_pipeline(s, y=1.95, h=2.35):
    """molecules -> fingerprints -> similarity -> 2-D layout."""
    steps = [("molecules", "a library\nyou care about"),
             ("fingerprint", "ECFP4, MHFP6,\ndescriptors"),
             ("similarity", "Tanimoto\nbetween all pairs"),
             ("layout", "t-SNE / UMAP /\nminimum spanning tree"),
             ("map", "one dot =\none molecule")]
    bw = 2.24
    gapx = (SW - 1.24 - 5 * bw) / 4
    for i, (t, sub) in enumerate(steps):
        x = 0.62 + i * (bw + gapx)
        acc = TEAL if i < 4 else ORANGE
        card(s, x, y, bw, h, fill=MINT if i < 4 else BLUSH)
        text(s, t, x, y + 0.22, bw, 0.32, size=13.5, color=acc, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, sub, x, y + 0.66, bw, 0.7, size=10.5, color=GRAY,
             align=PP_ALIGN.CENTER)
        if i < 4:
            arrow(s, x + bw + 0.04, y + h / 2, x + bw + gapx - 0.04, y + h / 2,
                  color=INK, w=1.8)
    text(s, "every choice in this chain changes the picture — the map is a property "
            "of your fingerprint, not of chemistry",
         0.62, y + h + 0.22, SW - 1.24, 0.3, size=12, color=RED, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h + 0.6


def fig_tmap(s, y=1.42, h=3.95):
    """Left: MST idea. Right: TMAP vs UMAP runtime (native log-ish chart)."""
    card(s, 0.62, y, 5.55, h, fill=CARD)
    panel_title(s, "TMAP: lay out a minimum spanning tree, not a scatter",
                0.85, y + 0.13, 5.2, size=11)
    # a little tree
    P = {0: (3.30, y + 1.05), 1: (2.35, y + 1.62), 2: (4.25, y + 1.58),
         3: (1.62, y + 2.30), 4: (2.95, y + 2.42), 5: (4.98, y + 2.20),
         6: (3.95, y + 2.72), 7: (1.30, y + 1.32), 8: (5.10, y + 1.15),
         9: (2.30, y + 3.00), 10: (4.60, y + 3.02)}
    E = [(0, 1), (0, 2), (1, 3), (1, 4), (2, 5), (2, 6), (1, 7), (2, 8),
         (3, 9), (6, 10)]
    for a, b in E:
        line(s, P[a][0], P[a][1], P[b][0], P[b][1], color=LGRAY, w=2.0)
    cols = {0: ORANGE, 1: TEAL, 2: TEAL, 3: PURPLE, 4: TEAL_L, 5: GREEN,
            6: TEAL_L, 7: PURPLE, 8: GREEN, 9: PURPLE, 10: TEAL_L}
    for k, (px, py) in P.items():
        dot(s, px, py, d=0.20 if k == 0 else 0.16, fill=cols[k])
    text(s, "branches = series of analogues · no overplotting at 10⁶ molecules",
         0.85, y + h - 0.34, 5.2, 0.3, size=11, color=INK,
         align=PP_ALIGN.CENTER)
    # right chart — log scale, or TMAP is invisible next to UMAP
    cats = ["10 k", "100 k", "500 k", "1 M"]
    ch = col_chart(s, 6.75, y + 0.42, 5.95, h - 1.10, cats,
                   [("TMAP", [4.9, 33.5, 175.9, 355.0]),
                    ("UMAP", [21.0, 115.7, 3578.0, 41326.0])],
                   ylabel="wall-clock seconds  (log scale)",
                   colors=[TEAL, ORANGE], legend=True, gap=55, tick_size=10)
    log_value_axis(ch)
    text(s, "1 M molecules:  TMAP 6 min   vs   UMAP 11.5 h   ( ~117× )",
         6.75, y + h - 0.40, 5.95, 0.32, size=12.5, color=RED, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_tsne_warnings(s, y=1.48, h=3.75):
    """Three things a neighbour embedding does NOT tell you."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    heads = ["gap width means nothing",
             "cluster size means nothing",
             "clusters appear in pure noise"]
    for x, hd in zip(xs, heads):
        card(s, x, y, W, h, fill=BLUSH)
        text(s, "✗  " + hd, x, y + 0.14, W, 0.3, size=12, color=RED, bold=True,
             align=PP_ALIGN.CENTER)
    # 1 gaps
    x = xs[0]
    blob(s, x + 1.05, y + 1.55, 0.55, 0.45, fill=WHITE, edge=TEAL)
    blob(s, x + 2.85, y + 1.55, 0.55, 0.45, fill=WHITE, edge=ORANGE)
    scatter_cloud(s, x + 1.05, y + 1.55, 0.40, 0.32, 9, 5, fill=TEAL, d=0.12)
    scatter_cloud(s, x + 2.85, y + 1.55, 0.40, 0.32, 9, 6, fill=ORANGE, d=0.12)
    line(s, x + 1.62, y + 1.55, x + 2.28, y + 1.55, color=GRAY, w=1.4, dash=DASH.DASH)
    text(s, "?", x + 1.80, y + 1.62, 0.3, 0.3, size=16, color=RED, bold=True)
    text(s, "perplexity changes it", x, y + h - 0.75, W, 0.3, size=10.5,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    # 2 sizes
    x = xs[1]
    blob(s, x + 1.15, y + 1.60, 0.72, 0.58, fill=WHITE, edge=TEAL)
    scatter_cloud(s, x + 1.15, y + 1.60, 0.54, 0.42, 7, 8, fill=TEAL, d=0.12)
    blob(s, x + 2.85, y + 1.60, 0.62, 0.50, fill=WHITE, edge=ORANGE)
    scatter_cloud(s, x + 2.85, y + 1.60, 0.46, 0.36, 24, 9, fill=ORANGE, d=0.09)
    text(s, "7 molecules", x + 0.28, y + 2.30, 1.7, 0.26, size=10.5, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "240 molecules", x + 1.98, y + 2.30, 1.7, 0.26, size=10.5, color=ORANGE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "drawn almost the same size", x, y + h - 0.68, W, 0.3, size=10.5,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    # 3 noise
    x = xs[2]
    for (bx, by, sd) in [(x + 1.05, y + 1.35, 31), (x + 2.75, y + 1.30, 32),
                         (x + 1.95, y + 2.35, 33)]:
        blob(s, bx, by, 0.50, 0.40, fill=WHITE, edge=LGRAY)
        scatter_cloud(s, bx, by, 0.36, 0.28, 8, sd, fill=GRAY, d=0.11)
    text(s, "input was uniform random", x, y + h - 0.75, W, 0.3, size=10.5,
         color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_butina(s, y=1.45, h=3.9):
    """Sphere-exclusion clustering in three frames."""
    W = 3.86
    xs = [0.62, 0.62 + W + 0.32, 0.62 + 2 * (W + 0.32)]
    heads = ["1 · count neighbours within t",
             "2 · densest molecule claims its sphere",
             "3 · repeat on what is left"]
    base = [(1.15, 1.25), (1.62, 1.60), (0.98, 1.85), (1.50, 2.20), (1.95, 1.15),
            (2.85, 1.35), (3.20, 1.95), (2.70, 2.30), (3.35, 2.62),
            (2.05, 2.85), (1.10, 2.70), (3.05, 1.05)]
    claimed = {0, 1, 2, 3, 4}
    for i, (x, hd) in enumerate(zip(xs, heads)):
        card(s, x, y, W, h, fill=CARD)
        text(s, hd, x, y + 0.14, W, 0.3, size=11.5, color=TEAL, bold=True,
             align=PP_ALIGN.CENTER)
        if i == 1:
            blob(s, x + base[1][0], y + base[1][1], 0.80, 0.72, fill=WHITE, edge=ORANGE)
        if i == 2:
            blob(s, x + base[1][0], y + base[1][1], 0.80, 0.72,
                 fill=RGBColor(0xF0, 0xF1, 0xF3), edge=LGRAY)
            blob(s, x + base[6][0], y + base[6][1], 0.72, 0.66, fill=WHITE, edge=PURPLE)
        for k, (bx, by) in enumerate(base):
            if i == 0:
                col = GRAY
            elif i == 1:
                col = ORANGE if k == 1 else (ORANGE_D if k in claimed else GRAY)
            else:
                col = LGRAY if k in claimed else (PURPLE if k == 6 else GRAY)
            dot(s, x + bx, y + by, d=0.20 if (i == 1 and k == 1) or (i == 2 and k == 6)
                else 0.14, fill=col)
        if i == 0:
            blob(s, x + base[1][0], y + base[1][1], 0.80, 0.72, fill=None,
                 edge=LGRAY) if False else None
            text(s, "t = Tanimoto cut-off", x, y + h - 0.44, W, 0.3, size=10.5,
                 color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        elif i == 1:
            text(s, "cluster centre + members", x, y + h - 0.44, W, 0.3, size=10.5,
                 color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
        else:
            text(s, "greedy, deterministic, no k to choose", x, y + h - 0.44, W, 0.3,
                 size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_scaffolds(s, y=1.4, h=4.05):
    """Bemis-Murcko: molecule -> framework, and how few frameworks cover drugs."""
    card(s, 0.62, y, 5.30, h, fill=CARD)
    panel_title(s, "strip side chains → keep the framework", 0.85, y + 0.13, 5.0, size=11.5)
    mol_sketch(s, 1.85, y + 1.52, sc=1.5, tails=3, seed=4, color=INK, hub=ORANGE)
    arrow(s, 2.90, y + 1.52, 3.55, y + 1.52, color=INK, w=2.0)
    mol_sketch(s, 4.55, y + 1.52, sc=1.5, tails=0, seed=4, color=TEAL)
    text(s, "molecule", 1.05, y + 2.40, 1.6, 0.26, size=10.5, color=GRAY,
         align=PP_ALIGN.CENTER)
    text(s, "Bemis–Murcko\nframework", 3.75, y + 2.40, 1.6, 0.5, size=10.5,
         color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    card(s, 0.85, y + h - 1.02, 4.85, 0.82, fill=BLUSH)
    text(s, "scaffold split = hold out whole frameworks,\nnot random molecules",
         0.85, y + h - 0.94, 4.85, 0.66, size=12, color=RED, bold=True,
         align=PP_ALIGN.CENTER)
    # cumulative coverage
    xy_chart(s, 6.35, y + 0.20, 6.35, h - 0.95,
             [{"name": "cumulative % of drugs covered",
               "pts": [(0, 0), (5, 21), (10, 32), (20, 42), (32, 50), (60, 60),
                       (120, 70), (300, 82), (700, 92), (1179, 100)],
               "kind": "both", "color": TEAL, "size": 6, "width": 2.6}],
             xlabel="number of distinct frameworks",
             ylabel="% of the 5,120 drugs", xlim=(0, 1200), ylim=(0, 100))
    card(s, 8.42, y + 1.58, 4.02, 0.96, fill=BLUSH)
    text(s, "32 frameworks", 8.42, y + 1.66, 4.02, 0.34, size=15, color=ORANGE_D,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "cover half of all known drugs", 8.42, y + 2.03, 4.02, 0.3, size=11.5,
         color=INK, align=PP_ALIGN.CENTER)
    text(s, "1,179 distinct frameworks among 5,120 drugs", 6.35, y + h - 0.30,
         6.35, 0.3, size=11, color=INK, italic=True, align=PP_ALIGN.CENTER)
    return y + h


def fig_periodic_embedding(s, y=1.42, h=4.0):
    """word2vec on 3.3M abstracts -> element vectors that rebuild the periodic table."""
    card(s, 0.62, y, 4.55, h, fill=MINT)
    panel_title(s, "no chemistry was taught", 0.85, y + 0.13, 4.2, size=11.5)
    rows = ["“… LiFePO₄ cathodes show …”",
            "“… we anneal the CoSb₃ skutterudite …”",
            "“… Bi₂Te₃ thermoelectric performance …”"]
    for k, t in enumerate(rows):
        rbox(s, 0.88, y + 0.58 + k * 0.52, 4.05, 0.42, fill=WHITE, edge=LGRAY,
             label=t, size=9.5, label_color=INK, bold=False)
    text(s, "3.3 M materials-science abstracts\n1.5 M used · word2vec · 200-D",
         0.88, y + 2.28, 4.05, 0.62, size=11, color=INK, align=PP_ALIGN.CENTER)
    text(s, "the only signal is which words co-occur", 0.88, y + h - 0.52, 4.05, 0.5,
         size=10.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    chevron(s, 5.28, y + 1.48, 1.15, 0.62, "t-SNE", fill=INK, size=11)
    # right: element groups emerge
    card(s, 6.62, y, 6.10, h, fill=CARD)
    panel_title(s, "t-SNE of the 100 element vectors", 6.86, y + 0.13, 5.6, size=11.5)
    groups = [("alkali", "Li Na K Rb Cs", 7.90, y + 1.05, ORANGE),
              ("halogens", "F Cl Br I", 11.35, y + 0.98, PURPLE),
              ("transition metals", "Fe Co Ni Cu", 9.62, y + 1.98, TEAL),
              ("chalcogens", "O S Se Te", 7.95, y + 2.92, GREEN),
              ("noble gases", "He Ne Ar", 11.45, y + 2.88, GRAY)]
    for lbl, syms, gx, gy, col in groups:
        blob(s, gx, gy, 0.62, 0.34, fill=WHITE, edge=col)
        scatter_cloud(s, gx, gy, 0.40, 0.18, 5, int(gx * 53), fill=col, d=0.11)
        text(s, lbl, gx - 0.95, gy + 0.36, 1.9, 0.26, size=9.5, color=col,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, syms, gx - 0.95, gy + 0.58, 1.9, 0.26, size=9, color=GRAY,
             align=PP_ALIGN.CENTER)
    text(s, "the periodic table falls out of the text alone",
         6.86, y + h - 0.36, 5.6, 0.3, size=12.5, color=TEAL, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h


def fig_rxnfp(s, y=1.45, h=3.95):
    """Reaction class recovery: hand-crafted fp vs self-supervised vs fine-tuned."""
    col_chart(s, 2.30, y, 8.75, h - 0.55,
              ["hand-crafted\nfingerprint + 5-NN",
               "self-supervised\npre-training only + 5-NN",
               "fine-tuned\nBERT + 5-NN"],
              [("classification accuracy", [41.0, 81.9, 98.9])],
              ylabel="accuracy over 792 reaction classes  (%)", ylim=(0, 105),
              colors=[[GRAY, TEAL, ORANGE]], data_labels=True, gap=70,
              tick_size=11)
    card(s, 5.80, y + 0.06, 2.50, 0.44, fill=MINT)
    text(s, "zero labels used", 5.80, y + 0.11, 2.50, 0.34, size=12, color=TEAL,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, "2.6 M unlabelled reactions in · reaction classes come out",
         2.30, y + h - 0.36, 8.75, 0.32, size=13, color=INK, bold=True,
         align=PP_ALIGN.CENTER)
    return y + h
