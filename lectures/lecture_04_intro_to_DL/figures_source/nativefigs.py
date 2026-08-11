#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Every figure in the deck, drawn as NATIVE PowerPoint objects:
  · data plots  -> real charts (right-click > Edit Data works)
  · schematics  -> real autoshapes and connectors (click and drag works)
  · tables      -> real PowerPoint tables
No raster images anywhere.
"""
import math
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import XyChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.oxml.ns import qn

from deck_lib import (PEACH, ORANGE, ORANGE_D, TEAL, TEAL_L, GREEN, PURPLE,
                      INK, GRAY, LGRAY, CARD, WHITE, FONT,
                      text, node, rbox, line, card, _margins)

MK_NONE, MK_CIRCLE, MK_TRIANGLE, MK_SQUARE = -4142, 8, 3, 1


# ===================================================================
# chart helpers
# ===================================================================
def _style_axis(ax, title=None, size=9.5, gridlines=False, tick_labels=True,
                lo=None, hi=None):
    ax.has_major_gridlines = gridlines
    if gridlines:
        gl = ax.major_gridlines.format.line
        gl.color.rgb = RGBColor(0xEC, 0xEE, 0xF2); gl.width = Pt(0.75)
    # Hide the axis line itself: on data that straddles zero it is drawn through the
    # middle of the plot and cuts across the curves. Gridlines carry the scale instead.
    ax.format.line.fill.background()
    if tick_labels:
        # LOW keeps the numbers at the edge of the plot; the default (NEXT_TO_AXIS)
        # draws them at the zero crossing, on top of the data
        ax.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        ax.tick_labels.font.size = Pt(size)
        ax.tick_labels.font.color.rgb = GRAY
        ax.tick_labels.font.name = FONT
        ax.tick_labels.number_format = "General"
        ax.tick_labels.number_format_is_linked = False
    else:
        ax.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    if lo is not None:
        ax.minimum_scale = lo
    if hi is not None:
        ax.maximum_scale = hi
    if title:
        ax.has_title = True
        ax.axis_title.text_frame.text = title
        r = ax.axis_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(size + 1); r.font.name = FONT
        r.font.color.rgb = GRAY; r.font.bold = False
    else:
        ax.has_title = False


def _pad_series(series):
    """python-pptx stacks XY series in the same spreadsheet columns. If a later
    series is longer than the first, some renderers drop the earlier series
    entirely. Padding every series to the same length with a repeat of its own
    last point is visually invisible and sidesteps the problem."""
    n = max((len(sp["pts"]) for sp in series), default=0)
    out = []
    for sp in series:
        q = dict(sp)
        pts = list(sp["pts"])
        if pts and len(pts) < n:
            pts = pts + [pts[-1]] * (n - len(pts))
        q["pts"] = pts
        out.append(q)
    return out


def xy_chart(slide, x, y, w, h, series, xlabel=None, ylabel=None,
             xlim=(None, None), ylim=(None, None), legend=False,
             gridlines=True, tick_size=9.5, tick_labels=True):
    """series: list of dicts {name, pts:[(x,y)...], kind:'markers'|'line'|'both',
    color, size, width, dash}."""
    series = _pad_series(series)
    cd = XyChartData()
    for sp in series:
        s = cd.add_series(sp.get("name", ""))
        for px, py in sp["pts"]:
            s.add_data_point(px, py)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_LINES,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.size = Pt(tick_size); ch.font.name = FONT
    for sp, ser in zip(series, ch.series):
        kind = sp.get("kind", "markers")
        col = sp.get("color", TEAL)
        if kind in ("markers", "both"):
            ser.marker.style = sp.get("marker", MK_CIRCLE)
            ser.marker.size = sp.get("size", 5)
            ser.marker.format.fill.solid()
            ser.marker.format.fill.fore_color.rgb = col
            ser.marker.format.line.color.rgb = sp.get("edge", col)
            ser.marker.format.line.width = Pt(0.75)
        else:
            ser.marker.style = MK_NONE
        if kind in ("line", "both"):
            lf = ser.format.line
            lf.color.rgb = col
            lf.width = Pt(sp.get("width", 2.25))
            if sp.get("dash"):
                lf.dash_style = sp["dash"]
        else:
            ser.format.line.fill.background()
        ser.smooth = sp.get("smooth", False)
    _style_axis(ch.value_axis, ylabel, tick_size, gridlines, tick_labels,
                ylim[0], ylim[1])
    _style_axis(ch.category_axis, xlabel, tick_size, gridlines, tick_labels,
                xlim[0], xlim[1])
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(tick_size)
        ch.legend.font.name = FONT
        ch.legend.font.color.rgb = INK
    # remove the chart's own outline so it sits cleanly on the white slide
    _no_chart_border(ch)
    return ch


def _no_chart_border(ch):
    from pptx.oxml.xmlchemy import OxmlElement as _OE
    cs = ch._chartSpace
    spPr = cs.find(qn('c:spPr'))
    if spPr is None:
        spPr = _OE('c:spPr')
        # c:spPr must follow c:chart inside c:chartSpace
        chart_el = cs.find(qn('c:chart'))
        if chart_el is not None:
            chart_el.addnext(spPr)
        else:
            cs.append(spPr)
    for tag in ('a:noFill', 'a:ln'):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    nf = _OE('a:noFill'); spPr.append(nf)
    ln = _OE('a:ln'); ln.append(_OE('a:noFill')); spPr.append(ln)


def col_chart(slide, x, y, w, h, cats, series, ylabel=None, ylim=(None, None),
              colors=None, legend=False, data_labels=False, gap=60,
              overlap=None, tick_size=9.5):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.size = Pt(tick_size); ch.font.name = FONT
    plot = ch.plots[0]
    plot.gap_width = gap
    if overlap is not None:
        plot.overlap = overlap
    if data_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(tick_size + 0.5); dl.font.bold = True
        dl.font.color.rgb = INK; dl.font.name = FONT
        dl.number_format = '0"%"'; dl.number_format_is_linked = False
    if colors:
        for ser, cols in zip(ch.series, colors):
            if isinstance(cols, list):     # per-point colouring
                for i, c in enumerate(cols):
                    pt = ser.points[i]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = c
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = cols
            ser.format.line.fill.background()
    _style_axis(ch.value_axis, ylabel, tick_size, True, True, ylim[0], ylim[1])
    _style_axis(ch.category_axis, None, tick_size, False, True)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(tick_size); ch.legend.font.name = FONT
        ch.legend.font.color.rgb = INK
    return ch


def panel_title(slide, t, x, y, w, size=12):
    text(slide, t, x, y, w, 0.28, size=size, color=INK, bold=True)


def caption(slide, t, x, y, w, size=10, color=GRAY):
    text(slide, t, x, y, w, 0.3, size=size, color=color, italic=True)


# ===================================================================
# 1. TIMELINE  (shapes)
# ===================================================================
def fig_timeline(slide, D, x=0.55, y=1.55, w=12.25, h=4.1):
    ev = D["timeline"]
    n = len(ev)
    axis_y = y + h * 0.50
    line(slide, x, axis_y, x + w, axis_y, color=GRAY, w=2.0)
    step = w / n
    for i, e in enumerate(ev):
        cx = x + step * (i + 0.5)
        up = (i % 2 == 1)
        col = ORANGE if int(e["year"]) >= 2012 else TEAL
        line(slide, cx, axis_y, cx, axis_y + (-0.28 if up else 0.28), color=col, w=1.8)
        d = 0.17
        node(slide, cx, axis_y, d, fill=col, edge=WHITE, edge_w=1.6)
        # year on the opposite side of the axis from the label
        ty = axis_y + (0.10 if up else -0.36)
        text(slide, e["year"], cx - step/2, ty, step, 0.26, size=11, color=col,
             bold=True, align=PP_ALIGN.CENTER)
        hy = axis_y - 0.72 if up else axis_y + 0.44
        text(slide, e["head"], cx - step/2, hy, step, 0.28, size=11.5, color=INK,
             bold=True, align=PP_ALIGN.CENTER)
        sy = axis_y - 1.22 if up else axis_y + 0.74
        text(slide, e["sub"].replace("|", "\n").split("\n")[0],
             cx - step/2, sy, step, 0.25, size=9.5, color=GRAY, align=PP_ALIGN.CENTER)
        text(slide, e["sub"].replace("|", "\n").split("\n")[1],
             cx - step/2, sy + 0.22, step, 0.25, size=9.5, color=GRAY,
             align=PP_ALIGN.CENTER)
    caption(slide, "The core ideas were in place by the 1980s — what changed was data, "
                   "GPUs and software.   (spacing not to scale)",
            x + 0.05, y + h - 0.18, w)


# ===================================================================
# 2. LINEAR FIT  (two native charts)
# ===================================================================
def fig_linear_fit(slide, D, y=1.45, h=3.15):
    d = D["linear_fit"]
    pts = list(zip(d["x"], d["y"]))
    w_, b_ = d["w"], d["b"]
    xs = [0.0, 5.0]
    panel_title(slide, f"Linear model:  ŷ = {w_:.2f}·x + {b_:.2f}", 0.85, y - 0.32, 6.0)
    xy_chart(slide, 0.75, y, 5.9, h,
             [{"name": "data", "pts": pts, "kind": "markers", "color": TEAL, "size": 6},
              {"name": "fit", "pts": [(v, w_*v + b_) for v in xs], "kind": "line",
               "color": ORANGE, "width": 2.5}],
             xlabel="feature  x", ylabel="target  y", xlim=(0, 5))
    imin = min(range(len(d["loss"])), key=lambda i: d["loss"][i])
    panel_title(slide, "Training = minimizing the loss over w", 7.05, y - 0.32, 6.0)
    xy_chart(slide, 6.95, y, 5.65, h,
             [{"name": "loss", "pts": list(zip(d["ws"], d["loss"])), "kind": "line",
               "color": TEAL, "width": 2.5},
              {"name": "minimum", "pts": [(d["ws"][imin], d["loss"][imin])],
               "kind": "markers", "color": ORANGE, "size": 9}],
             xlabel="weight  w", ylabel="loss  (MSE)")
    text(slide, f"one weight  →  one minimum  (w = {w_:.2f})", 7.05, y + h + 0.04,
         5.5, 0.28, size=10, color=ORANGE_D, italic=True)


# ===================================================================
# 3. XOR: linear vs one hidden layer  (charts + boundary polylines)
# ===================================================================
def _boundary_series(segs, color=ORANGE, width=2.5):
    out = []
    for s in segs:
        out.append({"name": "boundary", "pts": [(p[0], p[1]) for p in s],
                    "kind": "line", "color": color, "width": width, "smooth": True})
    return out


def fig_xor(slide, D, y=1.45, h=3.25):
    d = D["xor_2d"]
    c0 = [(p[0], p[1]) for p in d["class0"]]
    c1 = [(p[0], p[1]) for p in d["class1"]]
    base = [{"name": "class 0", "pts": c0, "kind": "markers", "color": TEAL, "size": 5},
            {"name": "class 1", "pts": c1, "kind": "markers", "color": ORANGE_D,
             "marker": MK_TRIANGLE, "size": 6}]
    panel_title(slide, f"Linear model — {d['acc_linear']:.0f}% accuracy",
                0.85, y - 0.32, 6.0)
    xy_chart(slide, 0.75, y, 5.9, h,
             base + _boundary_series(d["boundary_linear"]),
             xlim=(-2.6, 2.6), ylim=(-2.6, 2.6), tick_labels=False, legend=True)
    caption(slide, "no straight line can separate these classes", 0.85, y + h + 0.02, 5.8)
    panel_title(slide, f"One hidden layer (2→8→1) — {d['acc_mlp']:.0f}% accuracy",
                7.05, y - 0.32, 6.0)
    xy_chart(slide, 6.95, y, 5.65, h,
             base + _boundary_series(d["boundary_mlp"]),
             xlim=(-2.6, 2.6), ylim=(-2.6, 2.6), tick_labels=False)
    caption(slide, "a curved boundary solves it", 7.05, y + h + 0.02, 5.5)


# ===================================================================
# 4. ACTIVATIONS  (three native line charts)
# ===================================================================
def fig_activations(slide, D, y=1.65, h=2.75):
    d = D["activations"]
    z = d["z"]
    specs = [("ReLU   max(0, z)", "relu", ORANGE),
             ("tanh(z)", "tanh", TEAL),
             ("sigmoid   1/(1+e⁻ᶻ)", "sigmoid", PURPLE)]
    wpanel, gapx = 3.85, 0.30
    x0 = (13.3333 - (3*wpanel + 2*gapx)) / 2
    for i, (title, key, col) in enumerate(specs):
        xx = x0 + i * (wpanel + gapx)
        panel_title(slide, title, xx + 0.10, y - 0.32, wpanel)
        xy_chart(slide, xx, y, wpanel, h,
                 [{"name": key, "pts": list(zip(z, d[key])), "kind": "line",
                   "color": col, "width": 2.75}],
                 xlabel="z", xlim=(-3, 3), tick_size=9)


# ===================================================================
# 5. PARITY TABLE  (native table)
# ===================================================================
def fig_parity_table(slide, D, x=2.15, y=1.35, w=9.0):
    rows = D["parity_table"]
    heads = ["molecule", "nH", "nC", "nN", "nO", "electrons", "parity"]
    colw = [1.75, 1.0, 1.0, 1.0, 1.0, 1.75, 1.5]
    scale = w / sum(colw)
    colw = [c * scale for c in colw]
    n = len(rows) + 1
    h = 0.36 + 0.30 * len(rows)
    tb = slide.shapes.add_table(n, len(heads), Inches(x), Inches(y),
                                Inches(w), Inches(h)).table
    for i, cw in enumerate(colw):
        tb.columns[i].width = Inches(cw)
    for j, ht in enumerate(heads):
        c = tb.cell(0, j); c.text = ht
        c.fill.solid(); c.fill.fore_color.rgb = TEAL
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    pretty = {"CH4": "CH₄", "H2O": "H₂O", "NH3": "NH₃", "NO": "NO",
              "CH3": "•CH₃", "C6H6": "C₆H₆"}
    for i, rw in enumerate(rows):
        vals = ([pretty.get(rw["name"], rw["name"])] +
                [str(v) for v in rw["counts"]] +
                [str(rw["electrons"]), rw["parity"]])
        for j, v in enumerate(vals):
            c = tb.cell(i + 1, j); c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else WHITE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.runs[0]; r.font.size = Pt(12); r.font.name = FONT
            odd = (rw["parity"] == "odd")
            r.font.bold = (j == 0) or (j == 6)
            r.font.color.rgb = (ORANGE_D if (j == 6 and odd) else
                                (TEAL if j == 0 else INK))
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(slide, "electrons = 1·nH + 6·nC + 7·nN + 8·nO", x, y + h + 0.10, w/2, 0.3,
         size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "→   parity = (nH + nN) mod 2", x + w/2, y + h + 0.10, w/2, 0.3,
         size=13, color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    return y + h + 0.46


# ===================================================================
# 6. NETWORK DIAGRAMS  (shapes)
# ===================================================================
def draw_network(slide, W1, W2, labels, x, y, w, h, title=None,
                 dropout=None, node_d=0.30, in_d=0.44, show_legend=False):
    D_in, H = len(W1), len(W1[0])
    xs = [x + 0.10*w, x + 0.52*w, x + 0.93*w]
    wmax = max(max(abs(v) for row in W1 for v in row),
               max(abs(r[0]) for r in W2)) or 1.0
    pad_top = 0.42                      # room for the inputs/hidden/output captions
    top, bot = y + pad_top + 0.10, y + h - 0.16
    y_in = [top + (bot-top) * (i/(D_in-1)) for i in range(D_in)]
    y_hid = [y + pad_top + (h - pad_top - 0.12) * (j/(H-1)) for j in range(H)]
    y_out = y + pad_top + (h - pad_top)/2

    def edge(x1, y1, x2, y2, wt, alive=True):
        col = TEAL if wt >= 0 else ORANGE
        lw = 0.4 + 3.0*abs(wt)/wmax
        if not alive:
            col = RGBColor(0xE6, 0xE9, 0xEF); lw = 0.6
        line(slide, x1, y1, x2, y2, color=col, w=lw)

    for i in range(D_in):
        for j in range(H):
            edge(xs[0]+in_d/2, y_in[i], xs[1]-node_d/2, y_hid[j], W1[i][j],
                 True if dropout is None else bool(dropout[j]))
    for j in range(H):
        edge(xs[1]+node_d/2, y_hid[j], xs[2]-0.30, y_out, W2[j][0],
             True if dropout is None else bool(dropout[j]))
    for i, lab in enumerate(labels):
        node(slide, xs[0], y_in[i], in_d, fill=WHITE, edge=INK, label=lab,
             size=10.5, label_color=INK, edge_w=1.3)
    for j in range(H):
        dead = (dropout is not None and not dropout[j])
        node(slide, xs[1], y_hid[j], node_d,
             fill=WHITE if dead else TEAL, edge=LGRAY if dead else TEAL, edge_w=1.3)
        if dead:
            line(slide, xs[1]-0.09, y_hid[j]-0.09, xs[1]+0.09, y_hid[j]+0.09,
                 color=GRAY, w=1.4)
            line(slide, xs[1]-0.09, y_hid[j]+0.09, xs[1]+0.09, y_hid[j]-0.09,
                 color=GRAY, w=1.4)
    node(slide, xs[2], y_out, 0.44, fill=ORANGE, edge=ORANGE_D, label="ŷ",
         size=12, label_color=WHITE, edge_w=1.3)
    text(slide, "inputs", xs[0]-0.6, y + 0.04, 1.2, 0.25, size=9.5, color=GRAY,
         align=PP_ALIGN.CENTER)
    text(slide, f"hidden ({H})", xs[1]-0.8, y + 0.04, 1.6, 0.25, size=9.5,
         color=GRAY, align=PP_ALIGN.CENTER)
    text(slide, "output", xs[2]-0.6, y + 0.04, 1.2, 0.25, size=9.5, color=GRAY,
         align=PP_ALIGN.CENTER)
    if title:
        panel_title(slide, title, x + 0.05, y - 0.36, w)


def fig_net_init(slide, D, x=0.55, y=1.85, w=6.9, h=3.1):
    pn = D["parity_net"]
    draw_network(slide, pn["W1_init"], pn["W2_init"], pn["labels"], x, y, w, h,
                 title="Before training — the weights are random")
    return y + h


def fig_net_before_after(slide, D, y=1.95, h=3.0):
    pn = D["parity_net"]
    draw_network(slide, pn["W1_init"], pn["W2_init"], pn["labels"],
                 0.45, y, 4.4, h, title="initial (random)")
    draw_network(slide, pn["W1"], pn["W2"], pn["labels"],
                 4.75, y, 4.4, h, title="trained")
    panel_title(slide, "training loss", 9.55, y - 0.46, 3.4)
    hist = pn["loss_hist"]
    xy_chart(slide, 9.45, y, 3.45, h,
             [{"name": "loss", "pts": [(e, l) for e, l in hist], "kind": "line",
               "color": TEAL, "width": 2.25}],
             xlabel="epoch", ylabel="loss", tick_size=8.5)
    return y + h


# ===================================================================
# 7. LOSS SURFACE  (editable contour schematic + descent path)
# ===================================================================
def fig_loss_contour(slide, D, x=1.05, y=1.50, w=6.6, h=3.55):
    """Concentric ovals stand in for contour levels; the real descent path is drawn
    on top. Everything is a shape, so it can be recoloured or moved."""
    cx, cy = x + w*0.56, y + h*0.54
    shades = [RGBColor(0xEA, 0xF1, 0xF7), RGBColor(0xD3, 0xE4, 0xEF),
              RGBColor(0xB4, 0xD0, 0xE4), RGBColor(0x8E, 0xB8, 0xD5),
              RGBColor(0x63, 0x9C, 0xC2), RGBColor(0x2F, 0x7C, 0xA6)]
    for k, col in enumerate(shades):
        f = 1.0 - k*0.155
        ow, oh = w*0.86*f, h*0.80*f
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-ow/2),
                                    Inches(cy-oh/2), Inches(ow), Inches(oh))
        ov.rotation = 24
        ov.fill.solid(); ov.fill.fore_color.rgb = col
        ov.line.color.rgb = WHITE; ov.line.width = Pt(0.75)
        ov.shadow.inherit = False
    # descent path, mapped from weight space into the panel
    path = D["loss_surface"]["path"]
    pxs = [p[0] for p in path]; pys = [p[1] for p in path]
    x0, x1 = min(pxs), max(pxs); y0, y1 = min(pys), max(pys)
    def to_slide(p):
        fx = (p[0]-x0)/(x1-x0 or 1); fy = (p[1]-y0)/(y1-y0 or 1)
        return (x + w*(0.16 + 0.62*fx), y + h*(0.80 - 0.60*fy))
    pts = [to_slide(p) for p in path[::4]]
    for a, b in zip(pts[:-1], pts[1:]):
        line(slide, a[0], a[1], b[0], b[1], color=ORANGE, w=2.4)
    for p in pts[::3]:
        node(slide, p[0], p[1], 0.085, fill=ORANGE, edge=ORANGE, edge_w=0.5)
    node(slide, pts[0][0], pts[0][1], 0.20, fill=ORANGE_D, edge=WHITE, edge_w=1.4)
    node(slide, pts[-1][0], pts[-1][1], 0.20, fill=WHITE, edge=ORANGE_D, edge_w=2.0)
    text(slide, "start", pts[0][0]-0.75, pts[0][1]-0.34, 1.0, 0.25, size=9.5,
         color=ORANGE_D, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "after 60 steps", pts[-1][0]-0.35, pts[-1][1]+0.16, 1.9, 0.25,
         size=9.5, color=ORANGE_D, bold=True)
    text(slide, "w₁ →", x + w*0.44, y + h + 0.02, 0.9, 0.25, size=11, color=GRAY,
         bold=True)
    text(slide, "↑ w₂", x - 0.42, y + h*0.42, 0.7, 0.25, size=11, color=GRAY, bold=True)
    caption(slide, "loss contours over two weights (darker = lower loss)",
            x + 0.05, y + h + 0.30, w + 1.4)


# ===================================================================
# 8. LEARNING RATE  (three native charts)
# ===================================================================
def fig_learning_rate(slide, D, y=1.75, h=2.85):
    d = D["learning_rate"]
    curve = list(zip(d["curve_w"], d["curve_loss"]))
    panels = [("small", "η too small", "creeps — wastes your compute"),
              ("good",  "η about right", "converges quickly"),
              ("large", "η too large", "overshoots and diverges")]
    wp, gp = 3.95, 0.30
    x0 = (13.3333 - (3*wp + 2*gp)) / 2
    for i, (key, ttl, sub) in enumerate(panels):
        xx = x0 + i*(wp + gp)
        dd = d[key]
        panel_title(slide, f"{ttl}   (η = {dd['lr']})", xx + 0.08, y - 0.34, wp)
        xy_chart(slide, xx, y, wp, h,
                 [{"name": "loss", "pts": curve, "kind": "line", "color": TEAL,
                   "width": 2.0},
                  {"name": "steps", "pts": list(zip(dd["w"], dd["loss"])),
                   "kind": "both", "color": ORANGE, "width": 1.4, "size": 5}],
                 xlabel="weight w", ylim=(0, 14), tick_size=9)
        caption(slide, sub, xx + 0.08, y + h + 0.03, wp)


# ===================================================================
# 9. OVERFITTING: network capacity  (three native charts)
# ===================================================================
def fig_overfit_capacity(slide, D, y=1.70, h=2.95):
    d = D["overfit_capacity"]
    pts = list(zip(d["x"], d["y"]))
    truth = list(zip(d["xq"], d["truth"]))
    wp, gp = 3.95, 0.30
    x0 = (13.3333 - (3*wp + 2*gp)) / 2
    for i, (key, ttl) in enumerate([("h1", f"{d['labels']['h1']} — underfits"),
                                    ("h4", f"{d['labels']['h4']} — about right"),
                                    ("h64", f"{d['labels']['h64']} — overfits")]):
        xx = x0 + i*(wp + gp)
        fit = list(zip(d["xq"], d[key]))
        panel_title(slide, ttl, xx + 0.08, y - 0.34, wp)
        xy_chart(slide, xx, y, wp, h,
                 [{"name": "truth", "pts": truth, "kind": "line", "color": LGRAY,
                   "width": 1.75, "dash": 2},
                  {"name": "model", "pts": fit, "kind": "line", "color": ORANGE,
                   "width": 2.4},
                  {"name": "training data", "pts": pts, "kind": "markers",
                   "color": TEAL, "size": 5}],
                 ylim=(-2.4, 2.4), tick_size=9,
                 legend=(i == 0))


# ===================================================================
# 10. OVERFITTING SIGNATURE  (native chart + annotation shapes)
# ===================================================================
def fig_overfit_curves(slide, D, x=3.45, y=1.55, w=6.4, h=3.15):
    d = D["overfit_curves"]
    tr = list(zip(d["epoch"], d["train"]))
    va = list(zip(d["epoch"], d["val"]))
    be, bv = d["best_epoch"], d["best_val"]
    panel_title(slide, "The overfitting signature", x + 0.08, y - 0.34, w)
    ymax = max(d["val"]) * 1.05
    xy_chart(slide, x, y, w, h,
             [{"name": "training loss", "pts": tr, "kind": "line", "color": TEAL,
               "width": 2.4},
              {"name": "validation loss", "pts": va, "kind": "line", "color": ORANGE,
               "width": 2.4},
              {"name": "stop here", "pts": [(be, 0.0), (be, ymax)], "kind": "line",
               "color": INK, "width": 1.4, "dash": 2}],
             xlabel="epoch", ylabel="loss", legend=True, tick_size=9)
    text(slide, f"validation loss turns\naround at epoch {be}", x + w + 0.14,
         y + h*0.30, 2.7, 0.6, size=11.5, color=ORANGE_D, bold=True)
    text(slide, "beyond this point the model\nis memorizing, not learning",
         x + w + 0.14, y + h*0.30 + 0.62, 2.7, 0.7, size=10, color=GRAY, italic=True)
    return y + h


# ===================================================================
# 11. REGULARIZATION TRIO  (chart + chart + shapes)
# ===================================================================
def fig_regularization(slide, D, y=1.85, h=2.55):
    d = D["overfit_curves"]
    wd = D["weight_decay"]
    pn = D["parity_net"]
    wp = 3.85; gp = 0.42
    x0 = 0.55

    # (a) early stopping
    panel_title(slide, "Early stopping", x0 + 0.06, y - 0.36, wp)
    be = d["best_epoch"]
    ymax = max(d["val"]) * 1.05
    xy_chart(slide, x0, y, wp, h,
             [{"name": "train", "pts": list(zip(d["epoch"], d["train"])),
               "kind": "line", "color": TEAL, "width": 2.0},
              {"name": "val", "pts": list(zip(d["epoch"], d["val"])),
               "kind": "line", "color": ORANGE, "width": 2.0},
              {"name": "stop", "pts": [(be, 0), (be, ymax)], "kind": "line",
               "color": INK, "width": 1.4, "dash": 2}],
             xlabel="epoch", ylabel="loss", tick_size=8.5)
    caption(slide, "keep the weights from the validation minimum",
            x0 + 0.06, y + h + 0.04, wp + 0.3, size=9.5)

    # (b) weight decay
    x1 = x0 + wp + gp
    panel_title(slide, "Weight decay (L2)", x1 + 0.06, y - 0.36, wp)
    cats = [f"{c:.1f}" for c in wd["centers"]]
    col_chart(slide, x1, y, wp, h, cats,
              [("no penalty", wd["no_penalty"]), ("with weight decay", wd["with_decay"])],
              ylabel="count", colors=[LGRAY, ORANGE], legend=True, gap=20, overlap=100,
              tick_size=7.5)
    caption(slide, "add λ·Σwᵢ² → the weights shrink toward zero",
            x1 + 0.06, y + h + 0.04, wp + 0.3, size=9.5)

    # (c) dropout
    x2 = x1 + wp + gp
    panel_title(slide, "Dropout", x2 + 0.06, y - 0.36, wp)
    draw_network(slide, pn["W1"], pn["W2"], pn["labels"], x2, y + 0.10, wp, h - 0.10,
                 dropout=[1, 0, 1, 1, 0], node_d=0.26, in_d=0.38)
    caption(slide, "switch units off at random each step",
            x2 + 0.06, y + h + 0.04, wp + 0.3, size=9.5)


# ===================================================================
# 12. WEIGHT-DECAY-ONLY / SEEDS BAR CHART
# ===================================================================
def fig_seeds(slide, D, x=0.85, y=1.60, w=7.3, h=3.5):
    d = D["parity_seeds"]
    accs = d["acc"]
    cols = [ORANGE if a > 90 else (TEAL_L if a > 75 else GRAY) for a in accs]
    panel_title(slide, "Same network, different starting weights", x + 0.06, y - 0.36, w)
    ch = col_chart(slide, x, y, w, h, [str(s) for s in d["seeds"]],
                   [("test accuracy", accs)], ylabel="test accuracy  (%)",
                   ylim=(0, 105), colors=[cols], data_labels=True, gap=55)
    # chance line as a second series would rescale; draw it as a shape instead
    frac = 1 - 50/105.0
    ly = y + 0.16 + (h - 0.72) * frac
    line(slide, x + 0.62, ly, x + w - 0.15, ly, color=INK, w=1.4, dash=2)
    # the 50% reference is explained in the caption rather than labelled inside the
    # plot, where it would sit on top of a bar
    caption(slide, "dashed line = 50%, i.e. chance", x + 0.06, y + h + 0.06, w - 0.4)
    return y + h


# ===================================================================
# 13. CONVOLUTION + POOLING  (shapes)
# ===================================================================
def fig_conv_pool(slide, D, y=1.60, h=3.2):
    # --- convolution, left
    cell = 0.36
    rows, cols = 5, 7
    gx, gy = 1.15, y + 0.34
    panel_title(slide, "Convolution: one shared filter, slid", 0.85, y - 0.30, 6.0)
    for r in range(rows):
        for c in range(cols):
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(gx + c*cell), Inches(gy + r*cell),
                                       Inches(cell*0.9), Inches(cell*0.9))
            b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF7)
            b.line.color.rgb = RGBColor(0xCF, 0xD5, 0xE0); b.line.width = Pt(0.75)
            b.shadow.inherit = False
    def window(r0, c0, col, wdt, dash=None):
        wr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(gx + c0*cell - 0.02),
                                    Inches(gy + r0*cell - 0.02),
                                    Inches(cell*3 - 0.02), Inches(cell*3 - 0.02))
        wr.fill.background()
        wr.line.color.rgb = col; wr.line.width = Pt(wdt)
        if dash: wr.line.dash_style = dash
        wr.shadow.inherit = False
    window(0, 0, RGBColor(0xE8, 0xC6, 0xA6), 1.25, 2)
    window(1, 1, RGBColor(0xE8, 0xC6, 0xA6), 1.25, 2)
    window(1, 2, ORANGE, 2.5)
    caption(slide, "the same weights are reused at every position",
            gx - 0.30, gy + rows*cell + 0.10, 4.4, size=9.5)
    # arrow to feature map
    ax0 = gx + cols*cell + 0.18
    amid = gy + rows*cell/2
    line(slide, ax0, amid, ax0 + 0.55, amid, color=GRAY, w=1.6, arrow=True)
    # feature map
    fcell = 0.34
    fx, fy = ax0 + 0.80, gy + (rows*cell - 3*fcell)/2
    for r in range(3):
        for c in range(5):
            hit = (r == 1 and c == 2)
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(fx + c*fcell), Inches(fy + r*fcell),
                                       Inches(fcell*0.9), Inches(fcell*0.9))
            b.fill.solid()
            b.fill.fore_color.rgb = ORANGE if hit else RGBColor(0xCF, 0xE0, 0xF2)
            b.line.color.rgb = ORANGE_D if hit else RGBColor(0xAE, 0xC4, 0xDE)
            b.line.width = Pt(0.75); b.shadow.inherit = False
    text(slide, "feature map", fx - 0.20, fy - 0.32, 2.2, 0.25, size=9.5, color=GRAY)


def fig_pooling(slide, D, x=7.55, y=1.94, cell=0.55):
    grid = D["pooling"]["input"]; out = D["pooling"]["output"]
    panel_title(slide, "Pooling: keep the strongest response", x - 0.30, y - 0.64, 5.6)
    for r in range(4):
        for c in range(4):
            v = grid[r][c]
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x + c*cell), Inches(y + r*cell),
                                       Inches(cell), Inches(cell))
            shade = 0.10 + 0.09*v
            b.fill.solid()
            b.fill.fore_color.rgb = RGBColor(int(255-shade*150), int(255-shade*90),
                                             int(255-shade*30))
            b.line.color.rgb = WHITE; b.line.width = Pt(1.0)
            b.shadow.inherit = False
            tf = b.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rr = p.add_run(); rr.text = str(v)
            rr.font.size = Pt(12); rr.font.color.rgb = INK; rr.font.name = FONT
    for (r0, c0) in [(0, 0), (0, 2), (2, 0), (2, 2)]:
        wr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x + c0*cell), Inches(y + r0*cell),
                                    Inches(cell*2), Inches(cell*2))
        wr.fill.background(); wr.line.color.rgb = ORANGE; wr.line.width = Pt(2.25)
        wr.shadow.inherit = False
    ax0 = x + 4*cell + 0.12
    line(slide, ax0, y + 2*cell, ax0 + 0.5, y + 2*cell, color=GRAY, w=1.6, arrow=True)
    ox = ax0 + 0.72
    for r in range(2):
        for c in range(2):
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(ox + c*cell), Inches(y + cell + r*cell),
                                       Inches(cell), Inches(cell))
            b.fill.solid(); b.fill.fore_color.rgb = ORANGE
            b.line.color.rgb = WHITE; b.line.width = Pt(1.0)
            b.shadow.inherit = False
            tf = b.text_frame; _margins(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rr = p.add_run(); rr.text = str(out[r][c])
            rr.font.size = Pt(13); rr.font.bold = True
            rr.font.color.rgb = WHITE; rr.font.name = FONT
    text(slide, "max-pool 2×2", ox - 0.15, y + 3*cell + 0.10, 1.9, 0.25,
         size=9.5, color=GRAY)


# ===================================================================
# 14. MESSAGE PASSING  (shapes)
# ===================================================================
def fig_message_passing(slide, D, y=1.60, h=3.2):
    pos = {0: (0, 0), 1: (-1, .55), 2: (1, .55), 3: (-.75, -.85), 4: (.85, -.8),
           5: (-2.0, 1.05), 6: (2.0, 1.1)}
    edges = [(0, 1), (0, 2), (0, 3), (0, 4), (1, 5), (2, 6)]
    for panel, (px, rnd) in enumerate([(1.05, 1), (7.35, 2)]):
        sc = 1.25
        ox, oy = px + 2.4, y + h*0.52
        def P(k):
            return (ox + pos[k][0]*sc, oy - pos[k][1]*sc)
        panel_title(slide, f"Round {rnd}", px + 0.10, y - 0.30, 4.6)
        caption(slide, "reaches the direct neighbours" if rnd == 1
                else "information is now 2 bonds away",
                px + 0.10, y + h - 0.10, 4.8, size=9.5)
        for a, b in edges:
            xa, ya = P(a); xb, yb = P(b)
            line(slide, xa, ya, xb, yb, color=LGRAY, w=2.2)
        inc = [(1, 0), (2, 0), (3, 0), (4, 0)] if rnd == 1 else \
              [(5, 1), (6, 2), (1, 0), (2, 0), (3, 0), (4, 0)]
        for a, b in inc:
            xa, ya = P(a); xb, yb = P(b)
            dx, dy = xb-xa, yb-ya
            line(slide, xa+0.24*dx, ya+0.24*dy, xa+0.74*dx, ya+0.74*dy,
                 color=ORANGE, w=2.0, arrow=True)
        reached = {0, 1, 2, 3, 4} if rnd == 1 else set(pos)
        for k in pos:
            xk, yk = P(k)
            if k == 0:
                node(slide, xk, yk, 0.42, fill=ORANGE, edge=ORANGE_D, edge_w=1.5)
            else:
                node(slide, xk, yk, 0.34,
                     fill=TEAL if k in reached else WHITE, edge=TEAL, edge_w=1.5)


# ===================================================================
# 15. EQUIVARIANCE  (shapes)
# ===================================================================
def fig_equivariance(slide, D, y=1.55, h=3.0):
    import math as _m
    mol = [(0, 0), (.95, .35), (-.55, .85), (-.5, -.85)]
    frc = [(0, 0), (.55, .2), (-.3, .5), (-.28, -.5)]
    th = _m.radians(55)
    def rot(p):
        return (p[0]*_m.cos(th) - p[1]*_m.sin(th),
                p[0]*_m.sin(th) + p[1]*_m.cos(th))
    for panel, (px, do_rot) in enumerate([(1.35, False), (7.35, True)]):
        sc = 0.86
        ox, oy = px + 2.1, y + h*0.52
        M = [rot(p) if do_rot else p for p in mol]
        Fv = [rot(p) if do_rot else p for p in frc]
        panel_title(slide, "rotated molecule" if do_rot else "original molecule",
                    px + 0.05, y - 0.30, 4.4)
        for i in range(1, 4):
            line(slide, ox + M[0][0]*sc, oy - M[0][1]*sc,
                 ox + M[i][0]*sc, oy - M[i][1]*sc, color=LGRAY, w=2.4)
        for i, p in enumerate(M):
            node(slide, ox + p[0]*sc, oy - p[1]*sc, 0.30,
                 fill=ORANGE if i == 0 else TEAL, edge=WHITE, edge_w=1.4)
        for p, f in zip(M, Fv):
            if abs(f[0]) + abs(f[1]) < 1e-6:
                continue
            line(slide, ox + p[0]*sc, oy - p[1]*sc,
                 ox + (p[0]+f[0])*sc, oy - (p[1]+f[1])*sc,
                 color=PURPLE, w=2.0, arrow=True)
    text(slide, "energy  E = −152.3 kcal/mol", 1.40, y + h - 0.12, 4.4, 0.3,
         size=11.5, color=INK, bold=True)
    text(slide, "energy  E = −152.3 kcal/mol   (unchanged → invariant)",
         7.40, y + h - 0.12, 5.6, 0.3, size=11.5, color=TEAL, bold=True)
    text(slide, "the force vectors rotate with the molecule → equivariant",
         7.40, y + h + 0.14, 5.6, 0.3, size=11.5, color=PURPLE, bold=True)
    return y + h + 0.44


# ===================================================================
# 16. PARAMETER COUNT  (shapes + native table)
# ===================================================================
def fig_param_count(slide, D, y=1.55, h=3.0):
    pn = D["parity_net"]
    pc = D["param_count"]
    draw_network(slide, pn["W1"], pn["W2"], pn["labels"], 0.65, y + 0.30, 5.6, h - 0.30)
    rows = pc["rows"]
    x, tw = 7.1, 5.5
    n = len(rows) + 2
    th = 0.34 * n
    tb = slide.shapes.add_table(n, 3, Inches(x), Inches(y), Inches(tw),
                                Inches(th)).table
    for i, cw in enumerate([2.8, 1.4, 1.3]):
        tb.columns[i].width = Inches(cw)
    hdr = ["parameter block", "shape", "count"]
    for j, ht in enumerate(hdr):
        c = tb.cell(0, j); c.text = ht
        c.fill.solid(); c.fill.fore_color.rgb = TEAL
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(11.5); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.1)
    for i, (nm, shp, cnt) in enumerate(rows):
        for j, v in enumerate([nm, shp, str(cnt)]):
            c = tb.cell(i+1, j); c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else WHITE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.runs[0]; r.font.size = Pt(12); r.font.name = FONT
            r.font.color.rgb = GRAY if j == 1 else INK
            r.font.bold = (j == 2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.1)
    for j, v in enumerate(["total trainable parameters", "—", str(pc["total"])]):
        c = tb.cell(n-1, j); c.text = v
        c.fill.solid(); c.fill.fore_color.rgb = PEACH
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(13 if j == 2 else 12)
        r.font.bold = True; r.font.color.rgb = ORANGE_D; r.font.name = FONT
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.1)
    return max(y + th, y + h + 0.30)


# ===================================================================
# 17. AI / ML / DL NESTED BOXES  (shapes)
# ===================================================================
def fig_ai_ml_dl(slide, cx=10.15, cy=3.95):
    for (w, h, fill, lab, lc, ls, ly) in [
            (5.0, 3.4, CARD,  "ARTIFICIAL INTELLIGENCE", GRAY, 11, cy-1.57),
            (3.6, 2.35, PEACH, "MACHINE LEARNING", ORANGE_D, 11, cy-1.05),
            (2.15, 1.05, TEAL, "DEEP|LEARNING", WHITE, 13, cy-0.50)]:
        rbox(slide, cx-w/2, cy-h/2, w, h, fill=fill, edge=fill, label="")
        parts = lab.split("|")
        for k, part in enumerate(parts):
            text(slide, part, cx-w/2, ly + k*0.30, w, 0.3, size=ls, color=lc,
                 bold=True, align=PP_ALIGN.CENTER)
